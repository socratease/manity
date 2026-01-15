import asyncio
import io
import logging
import os
import re
import uuid
import argparse
from datetime import datetime
from enum import Enum
from email.message import EmailMessage
import smtplib
from pathlib import Path
from typing import List, Optional, Sequence

from fastapi import Body, Depends, FastAPI, File, HTTPException, Request, UploadFile, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, Response, JSONResponse
from pydantic import BaseModel, Field as PydanticField, field_validator
import httpx
from sqlalchemy import Column, ForeignKey, String, delete, event, func
from sqlalchemy.dialects.sqlite import JSON
from sqlalchemy.engine.url import make_url
from sqlalchemy.orm import selectinload
from sqlmodel import Field, Relationship, SQLModel, Session, create_engine, select

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

DEV_DEMO_SEED_ENV = "MANITY_ENABLE_DEMO_SEED"
ENVIRONMENT_ENV = "MANITY_ENV"
ADMIN_TOKEN_ENV = "MANITY_ADMIN_TOKEN"
PROTECTED_ENVIRONMENTS = {"prod", "production", "test", "testing"}

# Configure database path with persistent storage
# Default to persistent directory outside of application folder
DEFAULT_DEV_DB_PATH = "/home/c17420g/projects/manity-dev-data/portfolio.db"
DEFAULT_PROD_DB_PATH = "/home/c17420g/projects/manity-data/portfolio.db"
DEFAULT_DB_PATH = DEFAULT_DEV_DB_PATH
PERSISTENT_SQLITE_ROOTS = [
    Path("/var/data"),
    Path(DEFAULT_DEV_DB_PATH).parent,
    Path(DEFAULT_PROD_DB_PATH).parent,
]


def validate_sqlite_database_path(db_path: str | None) -> Path:
    if not db_path or db_path == ":memory:":
        logger.error(
            "DATABASE_URL cannot point to an in-memory SQLite database; configure a persistent path"
        )
        raise ValueError("DATABASE_URL must reference a persistent SQLite file")

    resolved_path = Path(db_path).expanduser()
    if not resolved_path.is_absolute():
        logger.error(
            "DATABASE_URL must use an absolute path for persistence (got %s)", resolved_path
        )
        raise ValueError("DATABASE_URL must be an absolute path for SQLite")

    if not any(resolved_path.is_relative_to(root) for root in PERSISTENT_SQLITE_ROOTS):
        logger.warning(
            "SQLite database path %s is outside known persistent mounts; data may not survive restarts",
            resolved_path,
        )

    resolved_path.parent.mkdir(parents=True, exist_ok=True)
    logger.info("Using SQLite database at %s", resolved_path)
    return resolved_path


def configure_sqlite_engine(engine):
    """Enable WAL mode so readers don't block writers and vice versa."""

    @event.listens_for(engine, "connect")
    def set_sqlite_pragma(dbapi_connection, connection_record):  # pragma: no cover - sqlite specific
        cursor = dbapi_connection.cursor()
        cursor.execute("PRAGMA journal_mode=WAL;")
        cursor.execute("PRAGMA synchronous=NORMAL;")
        cursor.execute("PRAGMA busy_timeout=30000;")
        cursor.execute("PRAGMA foreign_keys=ON;")
        cursor.close()


def create_engine_from_env(database_url: str | None = None):
    resolved_url = database_url or os.getenv("DATABASE_URL", f"sqlite:///{DEFAULT_DB_PATH}")
    url = make_url(resolved_url)

    if url.get_backend_name() == "sqlite":
        connect_args = {"check_same_thread": False, "timeout": 30}
        resolved_path = validate_sqlite_database_path(url.database)
        url = url.set(database=str(resolved_path))
    else:
        connect_args = {}
        logger.info(
            "Using database URL %s", url.render_as_string(hide_password=False)
        )

    engine = create_engine(url, connect_args=connect_args, pool_pre_ping=True)

    if url.get_backend_name() == "sqlite":
        configure_sqlite_engine(engine)

    return engine


engine = create_engine_from_env()
_PRAGMA_CACHE: dict[str, set[str]] = {}


def table_has_column(table_name: str, column_name: str) -> bool:
    cache_key = f"{table_name}:{column_name}"
    if cache_key in _PRAGMA_CACHE:
        return True

    with engine.connect() as connection:
        result = connection.exec_driver_sql(f"PRAGMA table_info({table_name})").all()
        for _, name, *_ in result:
            if name == column_name:
                _PRAGMA_CACHE[cache_key] = {column_name}
                return True
    return False


def ensure_column(table_name: str, column_definition: str) -> None:
    """
    Add a column to an existing table if it does not exist.

    SQLite does not support many ALTER operations, but adding nullable columns is safe.
    """
    column_name = column_definition.split()[0].strip('"')
    if table_has_column(table_name, column_name):
        return

    with engine.begin() as connection:
        connection.exec_driver_sql(f'ALTER TABLE "{table_name}" ADD COLUMN {column_definition}')
    _PRAGMA_CACHE[f"{table_name}:{column_name}"] = {column_name}


def column_has_unique_constraint(table_name: str, column_name: str) -> bool:
    """Check if a column has a UNIQUE constraint."""
    with engine.connect() as connection:
        # Get the CREATE TABLE statement
        result = connection.exec_driver_sql(
            f"SELECT sql FROM sqlite_master WHERE type='table' AND name='{table_name}'"
        ).fetchone()

        if not result:
            return False

        create_sql = result[0]
        if not create_sql:
            return False

        # Check for UNIQUE constraint on the column
        # Patterns: "column_name" UNIQUE, column_name UNIQUE, UNIQUE(column_name), etc.
        import re
        patterns = [
            rf'["`]?{column_name}["`]?\s+[^,]*?\bUNIQUE\b',  # column definition with UNIQUE
            rf'\bUNIQUE\s*\(\s*["`]?{column_name}["`]?\s*\)',  # UNIQUE(column)
        ]

        for pattern in patterns:
            if re.search(pattern, create_sql, re.IGNORECASE):
                return True

        return False


def migrate_add_unique_constraints(session: Session) -> None:
    """
    Add UNIQUE constraints to project.name and person.name.

    This migration:
    1. Checks if constraints already exist
    2. Handles duplicate names by appending numeric suffixes
    3. Recreates tables with UNIQUE constraints
    4. Preserves all relationships and foreign keys
    """
    migration_key = "add-unique-constraints-v1"
    if session.get(MigrationState, migration_key):
        logger.info("Migration %s already applied, skipping", migration_key)
        return

    logger.info("Running migration: %s", migration_key)

    # Check if constraints already exist
    project_has_unique = column_has_unique_constraint("project", "name")
    person_has_unique = column_has_unique_constraint("person", "name")

    if project_has_unique and person_has_unique:
        logger.info("UNIQUE constraints already exist, marking migration as complete")
        session.add(MigrationState(key=migration_key))
        session.commit()
        return

    with engine.begin() as connection:
        # Temporarily disable foreign key constraints during migration
        connection.exec_driver_sql("PRAGMA foreign_keys = OFF")

        # Migration for project.name
        if not project_has_unique:
            logger.info("Adding UNIQUE constraint to project.name")

            # First, resolve any duplicate project names
            result = connection.exec_driver_sql("""
                SELECT name, COUNT(*) as count
                FROM project
                GROUP BY LOWER(name)
                HAVING count > 1
            """).fetchall()

            if result:
                logger.warning("Found %d duplicate project names, resolving...", len(result))
                for name, count in result:
                    # Get all projects with this name (case-insensitive)
                    duplicates = connection.exec_driver_sql(
                        "SELECT id, name FROM project WHERE LOWER(name) = LOWER(?)",
                        [name]
                    ).fetchall()

                    # Keep the first one unchanged, rename the rest
                    for idx, (project_id, project_name) in enumerate(duplicates[1:], start=2):
                        new_name = f"{project_name} ({idx})"
                        logger.info("Renaming duplicate project '%s' -> '%s'", project_name, new_name)
                        connection.exec_driver_sql(
                            "UPDATE project SET name = ? WHERE id = ?",
                            [new_name, project_id]
                        )

            # Create new table with UNIQUE constraint
            connection.exec_driver_sql("""
                CREATE TABLE project_new (
                    id TEXT PRIMARY KEY,
                    name TEXT NOT NULL UNIQUE,
                    status TEXT NOT NULL,
                    priority TEXT NOT NULL,
                    progress INTEGER NOT NULL,
                    lastUpdate TEXT,
                    description TEXT NOT NULL,
                    executiveUpdate TEXT,
                    startDate TEXT,
                    targetDate TEXT,
                    stakeholders JSON
                )
            """)

            # Copy data
            connection.exec_driver_sql("""
                INSERT INTO project_new
                SELECT id, name, status, priority, progress, lastUpdate,
                       description, executiveUpdate, startDate, targetDate, stakeholders
                FROM project
            """)

            # Drop old table and rename new one
            connection.exec_driver_sql("DROP TABLE project")
            connection.exec_driver_sql("ALTER TABLE project_new RENAME TO project")

            # Recreate indexes
            connection.exec_driver_sql("CREATE INDEX ix_project_name ON project (name)")

            logger.info("Successfully added UNIQUE constraint to project.name")

        # Migration for person.name
        if not person_has_unique:
            logger.info("Adding UNIQUE constraint to person.name")

            # First, resolve any duplicate person names
            result = connection.exec_driver_sql("""
                SELECT name, COUNT(*) as count
                FROM person
                GROUP BY LOWER(name)
                HAVING count > 1
            """).fetchall()

            if result:
                logger.warning("Found %d duplicate person names, resolving...", len(result))
                for name, count in result:
                    # Get all people with this name (case-insensitive)
                    duplicates = connection.exec_driver_sql(
                        "SELECT id, name, team, email FROM person WHERE LOWER(name) = LOWER(?)",
                        [name]
                    ).fetchall()

                    # Merge logic: keep first one, update foreign keys from others
                    primary_id = duplicates[0][0]
                    duplicate_ids = [dup[0] for dup in duplicates[1:]]

                    logger.info("Merging %d duplicate entries for person '%s'", len(duplicate_ids), name)

                    # Update foreign keys to point to primary person
                    for dup_id in duplicate_ids:
                        connection.exec_driver_sql(
                            "UPDATE task SET assignee_id = ? WHERE assignee_id = ?",
                            [primary_id, dup_id]
                        )
                        connection.exec_driver_sql(
                            "UPDATE subtask SET assignee_id = ? WHERE assignee_id = ?",
                            [primary_id, dup_id]
                        )
                        connection.exec_driver_sql(
                            "UPDATE activity SET author_id = ? WHERE author_id = ?",
                            [primary_id, dup_id]
                        )
                        connection.exec_driver_sql(
                            "UPDATE projectpersonlink SET person_id = ? WHERE person_id = ?",
                            [primary_id, dup_id]
                        )

                        # Delete duplicate
                        connection.exec_driver_sql("DELETE FROM person WHERE id = ?", [dup_id])

            # Create new table with UNIQUE constraint
            connection.exec_driver_sql("""
                CREATE TABLE person_new (
                    id TEXT PRIMARY KEY,
                    name TEXT NOT NULL UNIQUE,
                    team TEXT NOT NULL,
                    email TEXT
                )
            """)

            # Copy data
            connection.exec_driver_sql("""
                INSERT INTO person_new
                SELECT id, name, team, email
                FROM person
            """)

            # Drop old table and rename new one
            connection.exec_driver_sql("DROP TABLE person")
            connection.exec_driver_sql("ALTER TABLE person_new RENAME TO person")

            # Recreate indexes
            connection.exec_driver_sql("CREATE INDEX ix_person_name ON person (name)")

            logger.info("Successfully added UNIQUE constraint to person.name")

        # Re-enable foreign key constraints
        connection.exec_driver_sql("PRAGMA foreign_keys = ON")

    # Mark migration as complete
    session.add(MigrationState(key=migration_key))
    session.commit()

    logger.info("Migration %s completed successfully", migration_key)


def migrate_remove_email_credentials(session: Session) -> None:
    """Remove email credential columns from the database.

    Email sending is now anonymous-only, so username/password columns are no longer needed.
    """
    migration_key = "remove-email-credentials-v1"
    if session.get(MigrationState, migration_key):
        logger.info("Migration %s already applied, skipping", migration_key)
        return

    logger.info("Running migration: %s", migration_key)

    with engine.begin() as connection:
        # Check if columns exist before trying to drop them
        result = connection.exec_driver_sql("PRAGMA table_info(emailsettings)")
        columns = [row[1] for row in result.fetchall()]

        if "username" not in columns and "password" not in columns:
            logger.info("Credential columns already removed, marking migration as complete")
            session.add(MigrationState(key=migration_key))
            session.commit()
            return

        # SQLite doesn't support DROP COLUMN in older versions, so we recreate the table
        connection.exec_driver_sql("PRAGMA foreign_keys = OFF")

        # Create new table without credential columns
        connection.exec_driver_sql("""
            CREATE TABLE IF NOT EXISTS emailsettings_new (
                id INTEGER PRIMARY KEY,
                smtp_server VARCHAR NOT NULL DEFAULT '',
                smtp_port INTEGER NOT NULL DEFAULT 587,
                use_tls BOOLEAN NOT NULL DEFAULT 1,
                from_address VARCHAR
            )
        """)

        # Copy data (excluding credentials)
        connection.exec_driver_sql("""
            INSERT INTO emailsettings_new (id, smtp_server, smtp_port, use_tls, from_address)
            SELECT id, smtp_server, smtp_port, use_tls, from_address FROM emailsettings
        """)

        # Swap tables
        connection.exec_driver_sql("DROP TABLE emailsettings")
        connection.exec_driver_sql("ALTER TABLE emailsettings_new RENAME TO emailsettings")

        connection.exec_driver_sql("PRAGMA foreign_keys = ON")

    # Mark migration as complete
    session.add(MigrationState(key=migration_key))
    session.commit()

    logger.info("Migration %s completed successfully - email credentials removed", migration_key)


def generate_id(prefix: str) -> str:
    return f"{prefix}-{uuid.uuid4().hex[:8]}"


def _normalize_env_value(value: str | None) -> str:
    return (value or "").strip().lower()


def current_environment() -> str:
    return _normalize_env_value(os.getenv(ENVIRONMENT_ENV, os.getenv("ENVIRONMENT")))


def is_dev_seeding_enabled() -> bool:
    environment = current_environment()
    if environment in PROTECTED_ENVIRONMENTS:
        logger.info("Skipping demo seeding because environment is set to %s", environment)
        return False

    flag_value = _normalize_env_value(os.getenv(DEV_DEMO_SEED_ENV))
    enabled = flag_value in {"1", "true", "yes", "on"}
    if not enabled:
        logger.info(
            "Demo project seeding disabled; set %s=1 to seed defaults in local development",
            DEV_DEMO_SEED_ENV,
        )
    return enabled


class Stakeholder(BaseModel):
    id: str | None = None
    name: str
    team: str = ""
    email: str | None = None


class PersonReference(SQLModel):
    id: Optional[str] = None
    name: Optional[str] = None
    team: Optional[str] = None
    email: Optional[str] = None


def normalize_stakeholders(stakeholders: Optional[List[Stakeholder | dict]]) -> list[dict]:
    normalized: list[dict] = []
    for stakeholder in stakeholders or []:
        if isinstance(stakeholder, Stakeholder):
            data = stakeholder.model_dump()
        elif isinstance(stakeholder, dict):
            data = {
                "id": stakeholder.get("id"),
                "name": stakeholder.get("name", ""),
                "team": stakeholder.get("team", ""),
                "email": stakeholder.get("email"),
            }
        else:  # pragma: no cover - defensive
            raise TypeError("Unsupported stakeholder type")

        data["name"] = (data.get("name") or "").strip()
        data["team"] = data.get("team") or ""
        data["email"] = (data.get("email") or None)
        normalized.append(data)
    return normalized


def serialize_person(person: Optional["Person"]) -> Optional[dict]:
    if person is None:
        return None
    return {
        "id": person.id,
        "name": person.name,
        "team": person.team,
        "email": person.email,
    }


def get_person_by_name(session: Session, name: str) -> "Person | None":
    if not name:
        return None

    normalized_name = name.strip()
    if not normalized_name:
        return None

    statement = select(Person).where(func.lower(Person.name) == normalized_name.lower())
    return session.exec(statement).first()


def get_person_by_email(session: Session, email: str | None) -> "Person | None":
    if not email:
        return None

    normalized_email = email.strip().lower()
    if not normalized_email:
        return None

    statement = select(Person).where(func.lower(Person.email) == normalized_email)
    return session.exec(statement).first()


class PersonIndex:
    def __init__(self, people: Sequence["Person"]):
        self.by_id: dict[str, Person] = {}
        self.by_name: dict[str, Person] = {}
        self.by_email: dict[str, Person] = {}

        for person in people:
            if person.id:
                self.by_id[person.id] = person
            if person.name:
                self.by_name[person.name.lower()] = person
            if person.email:
                self.by_email[person.email.lower()] = person

    def resolve(self, *, name: str | None = None, email: str | None = None, person_id: str | None = None) -> "Person | None":
        if person_id and person_id in self.by_id:
            return self.by_id[person_id]

        if email and (email.lower() in self.by_email):
            return self.by_email[email.lower()]

        if name and (name.lower() in self.by_name):
            return self.by_name[name.lower()]

        return None


def build_person_index(session: Session) -> PersonIndex:
    return PersonIndex(session.exec(select(Person)).all())


def _normalize_person_identity(name: str, email: str | None = None) -> tuple[str, str | None]:
    normalized_name = name.strip()
    normalized_email = email.strip().lower() if email else None
    return normalized_name, normalized_email


def _resolve_existing_person(
    session: Session,
    *,
    normalized_name: str,
    normalized_email: str | None = None,
    person_id: str | None = None,
) -> "Person | None":
    if person_id:
        person = session.get(Person, person_id)
        if person:
            return person

    person = get_person_by_email(session, normalized_email)
    if person:
        return person

    return get_person_by_name(session, normalized_name)


def upsert_person_from_payload(session: Session, payload: "PersonPayload") -> "Person":
    normalized_name, normalized_email = _normalize_person_identity(payload.name, payload.email)

    existing = _resolve_existing_person(
        session,
        normalized_name=normalized_name,
        normalized_email=normalized_email,
        person_id=payload.id,
    )

    if existing:
        existing.team = payload.team or existing.team
        existing.email = normalized_email or existing.email
        if normalized_name and existing.name.lower() != normalized_name.lower():
            conflict = get_person_by_name(session, normalized_name)
            if conflict is None or conflict.id == existing.id:
                existing.name = normalized_name
        session.add(existing)
        session.commit()
        session.refresh(existing)
        return existing

    person = Person(
        id=payload.id or generate_id("person"),
        name=normalized_name,
        team=payload.team,
        email=normalized_email,
    )
    session.add(person)
    session.commit()
    session.refresh(person)
    return person


def upsert_person_from_details(
    session: Session,
    name: str,
    team: str | None = None,
    email: str | None = None,
    person_id: str | None = None,
) -> "Person":
    normalized_name, normalized_email = _normalize_person_identity(name, email)
    normalized_team = team.strip() if team else ""

    existing = _resolve_existing_person(
        session,
        normalized_name=normalized_name,
        normalized_email=normalized_email,
        person_id=person_id,
    )

    if existing:
        if normalized_team:
            existing.team = normalized_team
        if email is not None:
            existing.email = normalized_email
        if normalized_name and existing.name.lower() != normalized_name.lower():
            conflict = get_person_by_name(session, normalized_name)
            if conflict is None or conflict.id == existing.id:
                existing.name = normalized_name
        session.add(existing)
        session.commit()
        session.refresh(existing)
        return existing

    person = Person(
        id=person_id or generate_id("person"),
        name=normalized_name,
        team=normalized_team,
        email=normalized_email,
    )
    session.add(person)
    session.commit()
    session.refresh(person)
    return person


def resolve_person_reference(session: Session, reference) -> "Person | None":
    """
    Accepts a variety of person representations (id dict, PersonPayload, Stakeholder, or name string)
    and returns a persisted Person instance, creating or updating as needed.
    """
    if reference is None:
        return None

    if isinstance(reference, Person):
        return reference

    if isinstance(reference, str):
        normalized = reference.strip()
        if not normalized:
            return None
        payload = PersonPayload(name=normalized, team="Contributor")
        return upsert_person_from_payload(session, payload)

    person_id = None
    name = None
    team = None
    email = None

    if isinstance(reference, Stakeholder):
        person_id = reference.id
        name = reference.name
        team = reference.team
    elif isinstance(reference, AssigneePayload):
        person_id = reference.id
        name = reference.name
        team = reference.team
    elif isinstance(reference, PersonPayload):
        person_id = reference.id
        name = reference.name
        team = reference.team
        email = reference.email
    elif isinstance(reference, PersonReference):
        person_id = reference.id
        name = reference.name
        team = reference.team
        email = reference.email
    elif isinstance(reference, dict):
        person_id = reference.get("id")
        name = reference.get("name")
        team = reference.get("team")
        email = reference.get("email")
    else:  # pragma: no cover - defensive
        return None

    normalized_name = (name or "").strip()
    normalized_team = (team or "").strip() or "Contributor"

    if person_id:
        person = session.get(Person, person_id)
        if person:
            if normalized_name:
                person.name = normalized_name
            person.team = normalized_team or person.team
            if email is not None:
                person.email = email
            session.add(person)
            session.commit()
            session.refresh(person)
            return person

        if not normalized_name:
            return None

        person = Person(
            id=person_id,
            name=normalized_name,
            team=normalized_team,
            email=email,
        )
        session.add(person)
        session.commit()
        session.refresh(person)
        return person

    if not normalized_name:
        return None

    payload = PersonPayload(name=normalized_name, team=normalized_team, email=email)
    return upsert_person_from_payload(session, payload)


class SubtaskBase(SQLModel):
    title: str
    status: str = "todo"
    dueDate: Optional[str] = None
    completedDate: Optional[str] = None
    assignee_id: Optional[str] = Field(
        default=None,
        sa_column=Column("assignee_id", String, ForeignKey("person.id", ondelete="SET NULL"), nullable=True),
        description="Person responsible for this subtask",
        alias="assigneeId",
    )


class Subtask(SubtaskBase, table=True):
    id: Optional[str] = Field(default=None, primary_key=True)
    task_id: Optional[str] = Field(default=None, foreign_key="task.id")
    assignee_id: Optional[str] = Field(default=None, foreign_key="person.id")
    task: "Task" = Relationship(back_populates="subtasks")
    assignee: Optional["Person"] = Relationship(sa_relationship_kwargs={"lazy": "joined"})


class TaskBase(SQLModel):
    title: str
    status: str = "todo"
    dueDate: Optional[str] = None
    completedDate: Optional[str] = None
    assignee_id: Optional[str] = Field(
        default=None,
        sa_column=Column("assignee_id", String, ForeignKey("person.id", ondelete="SET NULL"), nullable=True),
        description="Person responsible for this task",
        alias="assigneeId",
    )


class Task(TaskBase, table=True):
    id: Optional[str] = Field(default=None, primary_key=True)
    project_id: Optional[str] = Field(default=None, foreign_key="project.id")
    assignee_id: Optional[str] = Field(default=None, foreign_key="person.id")
    project: "Project" = Relationship(back_populates="plan")
    assignee: Optional["Person"] = Relationship(sa_relationship_kwargs={"lazy": "joined"})
    subtasks: list[Subtask] = Relationship(
        back_populates="task",
        sa_relationship_kwargs={"cascade": "all, delete-orphan"},
    )


class ActivityBase(SQLModel):
    date: str
    note: str
    author: Optional[str] = None
    author_id: Optional[str] = Field(
        default=None,
        sa_column=Column("author_id", String, ForeignKey("person.id", ondelete="SET NULL"), nullable=True),
        alias="authorId",
    )


class Activity(ActivityBase, table=True):
    id: Optional[str] = Field(default=None, primary_key=True)
    project_id: Optional[str] = Field(default=None, foreign_key="project.id")
    # Store task context as JSON for comments on tasks/subtasks
    task_context: Optional[str] = Field(default=None, sa_column=Column(String))
    project: "Project" = Relationship(back_populates="recentActivity")
    author_person: Optional["Person"] = Relationship(sa_relationship_kwargs={"lazy": "joined"})


class ProjectPersonLink(SQLModel, table=True):
    project_id: str = Field(
        sa_column=Column("project_id", String, ForeignKey("project.id", ondelete="CASCADE"), primary_key=True, nullable=False),
    )
    person_id: str = Field(
        sa_column=Column("person_id", String, ForeignKey("person.id", ondelete="CASCADE"), primary_key=True, nullable=False),
    )


class InitiativePersonLink(SQLModel, table=True):
    """Link table for many-to-many relationship between Initiatives and People (owners)."""
    initiative_id: str = Field(
        sa_column=Column("initiative_id", String, ForeignKey("initiative.id", ondelete="CASCADE"), primary_key=True, nullable=False),
    )
    person_id: str = Field(
        sa_column=Column("person_id", String, ForeignKey("person.id", ondelete="CASCADE"), primary_key=True, nullable=False),
    )


class InitiativeBase(SQLModel):
    """Base model for Initiative with common fields."""
    name: str = Field(sa_column=Column(String, unique=True, index=True))
    description: str = ""
    status: str = "planning"  # planning | active | on-hold | cancelled | completed
    priority: str = "medium"  # high | medium | low
    startDate: Optional[str] = None
    targetDate: Optional[str] = None


class Initiative(InitiativeBase, table=True):
    """Initiative entity - a meta-project that groups related projects."""
    id: Optional[str] = Field(default=None, primary_key=True)

    # Relationships
    projects: list["Project"] = Relationship(back_populates="initiative")
    owners: list["Person"] = Relationship(
        back_populates="owned_initiatives",
        link_model=InitiativePersonLink,
    )


class ProjectBase(SQLModel):
    name: str = Field(sa_column=Column(String, unique=True, index=True))
    status: str = "planning"
    priority: str = "medium"
    progress: int = 0
    lastUpdate: Optional[str] = None
    description: str = ""
    executiveUpdate: Optional[str] = None
    startDate: Optional[str] = None
    targetDate: Optional[str] = None
    stakeholders_legacy: List[Stakeholder] = Field(
        default_factory=list,
        sa_column=Column("stakeholders", JSON, nullable=True),
    )


class Project(ProjectBase, table=True):
    id: Optional[str] = Field(default=None, primary_key=True)
    initiative_id: Optional[str] = Field(
        default=None,
        sa_column=Column("initiative_id", String, ForeignKey("initiative.id", ondelete="SET NULL"), nullable=True),
    )
    plan: list[Task] = Relationship(
        back_populates="project",
        sa_relationship_kwargs={"cascade": "all, delete-orphan"},
    )
    recentActivity: list[Activity] = Relationship(
        back_populates="project",
        sa_relationship_kwargs={"cascade": "all, delete-orphan"},
    )
    stakeholders: list["Person"] = Relationship(
        back_populates="projects",
        link_model=ProjectPersonLink,
    )
    initiative: Optional["Initiative"] = Relationship(back_populates="projects")


class PersonBase(SQLModel):
    name: str = Field(sa_column=Column(String, unique=True, index=True))
    team: str = ""
    email: Optional[str] = None


class Person(PersonBase, table=True):
    id: Optional[str] = Field(default=None, primary_key=True)
    projects: list["Project"] = Relationship(
        back_populates="stakeholders",
        link_model=ProjectPersonLink,
    )
    owned_initiatives: list["Initiative"] = Relationship(
        back_populates="owners",
        link_model=InitiativePersonLink,
    )


class EmailSettings(SQLModel, table=True):
    id: Optional[int] = Field(default=1, primary_key=True)
    smtp_server: str = ""
    smtp_port: int = 587
    use_tls: bool = True
    from_address: Optional[str] = None


class AuditLog(SQLModel, table=True):
    """Audit log for tracking all actions and AI conversations"""
    id: Optional[int] = Field(default=None, primary_key=True)
    timestamp: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    action: str  # e.g., "create_project", "update_task", "llm_chat"
    entity_type: Optional[str] = None  # e.g., "project", "task", "person"
    entity_id: Optional[str] = None  # ID of the affected entity
    details: Optional[str] = Field(default=None, sa_column=Column(String))  # JSON string with additional details
    user_agent: Optional[str] = None  # Client user agent
    ip_address: Optional[str] = None  # Client IP address


class MigrationState(SQLModel, table=True):
    """Track lightweight migrations run in-application."""

    key: str = Field(primary_key=True)
    applied_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())


def log_action(
    session: Session,
    action: str,
    entity_type: str = None,
    entity_id: str = None,
    details: dict = None,
    request: Request = None
):
    """Log an action to the audit log"""
    import json

    log_entry = AuditLog(
        action=action,
        entity_type=entity_type,
        entity_id=entity_id,
        details=json.dumps(details) if details else None,
        user_agent=request.headers.get("user-agent") if request else None,
        ip_address=request.client.host if request and request.client else None
    )
    session.add(log_entry)
    session.commit()
    logger.info(f"Action logged: {action} on {entity_type}:{entity_id}")


def get_logged_in_user(request: Request | None) -> str | None:
    if not request:
        return None
    for header_name in ("x-logged-in-user", "x-user-name", "x-user"):
        header_value = request.headers.get(header_name)
        if header_value and header_value.strip():
            return header_value.strip()
    return None


def resolve_activity_author(
    session: Session,
    request: Request | None,
    fallback: str | None = None
) -> tuple[str, str | None]:
    name = get_logged_in_user(request) or fallback
    if name:
        author_person = resolve_person_reference(session, name)
        return author_person.name if author_person else name, author_person.id if author_person else None
    return "Unknown", None


class AssigneePayload(BaseModel):
    """Assignee information - can include id, name, or both"""
    id: Optional[str] = None
    name: Optional[str] = None
    team: Optional[str] = None


class TaskContextPayload(BaseModel):
    """Task context for comments on tasks/subtasks"""
    taskId: str
    subtaskId: Optional[str] = None
    taskTitle: str
    subtaskTitle: Optional[str] = None


class SubtaskPayload(SubtaskBase):
    id: Optional[str] = None
    assignee: Optional[AssigneePayload] = None


class TaskPayload(TaskBase):
    id: Optional[str] = None
    subtasks: Optional[List[SubtaskPayload]] = None
    assignee: Optional[AssigneePayload] = None


class ActivityPayload(ActivityBase):
    id: Optional[str] = None
    taskContext: Optional[TaskContextPayload] = None
    authorEmail: Optional[str] = None


class PersonPayload(PersonBase):
    id: Optional[str] = None


class ProjectPayload(SQLModel):
    name: str
    status: str = "planning"
    priority: str = "medium"
    progress: int = 0
    lastUpdate: Optional[str] = None
    description: str = ""
    executiveUpdate: Optional[str] = None
    startDate: Optional[str] = None
    targetDate: Optional[str] = None
    stakeholders: List[PersonReference] = Field(default_factory=list)
    id: Optional[str] = None
    plan: List[TaskPayload] = Field(default_factory=list)
    recentActivity: List[ActivityPayload] = Field(default_factory=list)

    @field_validator("name", mode="before")
    @classmethod
    def validate_name(cls, v):
        if v is None:
            raise ValueError("Project name cannot be null")
        if isinstance(v, str):
            v = v.strip()
            if not v:
                raise ValueError("Project name cannot be empty")
        return v


class InitiativePayload(SQLModel):
    """Payload for creating/updating an initiative."""
    id: Optional[str] = None
    name: str
    description: str = ""
    status: str = "planning"
    priority: str = "medium"
    startDate: Optional[str] = None
    targetDate: Optional[str] = None
    owners: List[PersonReference] = Field(default_factory=list)

    @field_validator("name", mode="before")
    @classmethod
    def validate_name(cls, v):
        if v is None:
            raise ValueError("Initiative name cannot be null")
        if isinstance(v, str):
            v = v.strip()
            if not v:
                raise ValueError("Initiative name cannot be empty")
        return v


class ImportPayload(BaseModel):
    projects: List[ProjectPayload]
    people: List[PersonPayload] = Field(default_factory=list)
    mode: str = "replace"


class EmailSettingsPayload(BaseModel):
    smtpServer: str
    smtpPort: int = 587
    useTLS: bool = True
    fromAddress: Optional[str] = None


class EmailSettingsResponse(BaseModel):
    smtpServer: str
    smtpPort: int
    fromAddress: Optional[str] = None
    useTLS: bool = True


class EmailSendPayload(BaseModel):
    recipients: List[str] | str
    cc: Optional[List[str] | str] = None
    bcc: Optional[List[str] | str] = None
    subject: str
    body: str
    # Inline SMTP settings (sent with each request, stored in browser)
    smtp_server: Optional[str] = None
    smtp_port: Optional[int] = 587
    from_address: Optional[str] = None
    use_tls: Optional[bool] = True


def get_email_settings(session: Session) -> EmailSettings:
    settings = session.get(EmailSettings, 1)
    if settings is None:
        settings = EmailSettings(id=1)
        session.add(settings)
        session.commit()
        session.refresh(settings)
    return settings


def serialize_email_settings(settings: EmailSettings) -> dict:
    return {
        "smtpServer": settings.smtp_server,
        "smtpPort": settings.smtp_port,
        "fromAddress": settings.from_address or "",
        "useTLS": settings.use_tls,
    }


def normalize_recipients(raw: Sequence[str] | str) -> list[str]:
    if isinstance(raw, str):
        candidates = [raw]
    else:
        candidates = list(raw)

    recipients: list[str] = []
    for candidate in candidates:
        if not candidate:
            continue
        if isinstance(candidate, str):
            parts = candidate.replace(";", ",").split(",")
            for part in parts:
                normalized = part.strip()
                if normalized:
                    recipients.append(normalized)
    return recipients


def dispatch_email(
    smtp_server: str,
    smtp_port: int,
    from_address: str,
    recipients: list[str],
    cc: list[str],
    bcc: list[str],
    subject: str,
    body: str,
    use_tls: bool = False
) -> dict:
    """
    Send an email via SMTP anonymously (no authentication).

    The server is expected to be a local or trusted SMTP relay.

    Returns a dict with 'sent_to' (list of successful recipients) and any 'refused' recipients.
    Raises ValueError for configuration issues, SMTPException for server errors.
    """
    if not smtp_server:
        raise ValueError("SMTP server is not configured. Please set the server address in settings.")
    if not from_address:
        raise ValueError("Sender address is not configured. Please set the From address in settings.")
    if not (recipients or cc or bcc):
        raise ValueError("At least one recipient is required")

    message = EmailMessage()
    message["Subject"] = subject
    message["From"] = from_address
    if recipients:
        message["To"] = ", ".join(recipients)
    if cc:
        message["Cc"] = ", ".join(cc)
    message.set_content(body)
    all_recipients = [*recipients, *cc, *bcc]

    try:
        with smtplib.SMTP(smtp_server, smtp_port, timeout=30) as smtp:
            # Check initial connection when supported
            if hasattr(smtp, "noop"):
                code, resp = smtp.noop()
                if code != 250:
                    raise smtplib.SMTPException(f"Server not ready: {code} {resp.decode()}")

            # Use STARTTLS only if explicitly requested
            if use_tls:
                try:
                    if hasattr(smtp, "starttls"):
                        result = smtp.starttls()
                        if isinstance(result, tuple):
                            code, resp = result
                            if code != 220:
                                logger.warning("STARTTLS failed with code %d, continuing without TLS", code)
                    else:
                        logger.info("Server does not support STARTTLS, sending without encryption")
                except smtplib.SMTPNotSupportedError:
                    logger.info("Server does not support STARTTLS, sending without encryption")

            # send_message returns dict of refused recipients (empty = all accepted)
            refused = smtp.send_message(message, to_addrs=all_recipients)

            # Verify all recipients were accepted
            if refused:
                refused_addrs = list(refused.keys())
                logger.warning("Some recipients refused: %s", refused_addrs)
                if len(refused_addrs) == len(all_recipients):
                    raise smtplib.SMTPRecipientsRefused(refused)

            # Verify message was queued by checking server response
            # Issue a NOOP after sending to confirm connection is still good
            if hasattr(smtp, "noop"):
                code, resp = smtp.noop()
                if code != 250:
                    logger.warning("Post-send NOOP returned %d: %s", code, resp.decode())

            successful = [r for r in all_recipients if r not in (refused or {})]
            logger.info("Email sent successfully to %d recipient(s): %s", len(successful), successful)

            return {
                "sent_to": successful,
                "refused": list(refused.keys()) if refused else []
            }

    except smtplib.SMTPRecipientsRefused as exc:
        logger.exception("All recipients refused")
        raise ValueError(f"All recipients refused by server")
    except smtplib.SMTPConnectError as exc:
        logger.exception("Failed to connect to SMTP server")
        raise ValueError(f"Could not connect to email server at {smtp_server}:{smtp_port}. Please check your settings.")
    except smtplib.SMTPException as exc:
        logger.exception("SMTP error sending email")
        raise exc
    except ConnectionRefusedError:
        logger.exception("Connection refused by SMTP server")
        raise ValueError(f"Connection refused by email server at {smtp_server}:{smtp_port}. Please verify the server is running.")
    except Exception as exc:
        logger.exception("Failed to send email")
        raise exc


class ChatRole(str, Enum):
    SYSTEM = "system"
    USER = "user"
    ASSISTANT = "assistant"
    TOOL = "tool"


class ToolCallFunction(BaseModel):
    name: str
    arguments: str


class ToolCall(BaseModel):
    id: str
    type: str = "function"
    function: ToolCallFunction


class ChatMessage(BaseModel):
    role: ChatRole
    content: Optional[str] = None
    tool_calls: Optional[List[ToolCall]] = None
    tool_call_id: Optional[str] = None
    name: Optional[str] = None  # For tool messages


class ChatProvider(str, Enum):
    OPENAI = "openai"
    AZURE_OPENAI = "azure"


class ToolDefinition(BaseModel):
    type: str = "function"
    function: dict


class ChatRequest(BaseModel):
    model: str = os.getenv("LLM_MODEL", "gpt-5.1")
    provider: Optional[ChatProvider] = None
    messages: List[ChatMessage] = PydanticField(..., min_length=1)
    response_format: Optional[dict] = None
    tools: Optional[List[ToolDefinition]] = None
    tool_choice: Optional[str] = None


FRONTEND_ORIGINS_ENV = "FRONTEND_ORIGINS"
FRONTEND_ORIGIN_REGEX_ENV = "FRONTEND_ORIGIN_REGEX"


def parse_origins(value: str | None) -> list[str]:
    return [origin.strip() for origin in (value or "").split(",") if origin.strip()]


def configured_frontend_origins() -> tuple[list[str], str | None]:
    origins = parse_origins(os.getenv(FRONTEND_ORIGINS_ENV))
    origin_regex = os.getenv(FRONTEND_ORIGIN_REGEX_ENV) or None

    if not origins and not origin_regex:
        # Default: allow all localhost/127.0.0.1 origins for local development
        # This covers common dev ports: 3000, 5173, 8113, 8114, etc.
        origin_regex = r"^https?://(localhost|127\.0\.0\.1|rn000224)(:\d+)?$"

    return origins, origin_regex


app = FastAPI(title="Manity Portfolio API")

# CORS configuration
# Note: allow_credentials=True with allow_origins=["*"] violates CORS spec.
# When credentials are needed, the server must echo the specific origin.
# Setting allow_credentials=False allows the wildcard origin to work properly.
# If cookie/auth credentials are needed, specify explicit origins in CORS_ORIGINS env var.

# Echo the requesting origin when possible so browsers see a concrete value
# instead of "*". When an explicit wildcard is requested, fall back to a
# permissive regex that still mirrors the Origin header while keeping
# allow_credentials disabled.
allowed_origins, allowed_origin_regex = configured_frontend_origins()

logger.info("CORS allowed_origins=%s", allowed_origins)
logger.info("CORS allowed_origin_regex=%s", allowed_origin_regex)

if "*" in allowed_origins and not allowed_origin_regex:
    allowed_origin_regex = ".*"
    allowed_origins = []

logger.info("CORS allowed_origins=%s", allowed_origins)
logger.info("CORS allowed_origin_regex=%s", allowed_origin_regex)

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_origin_regex=allowed_origin_regex,
    allow_credentials=False,   # ✅ no cookies, no sessions
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
)

def create_db_and_tables() -> None:
    SQLModel.metadata.create_all(engine)
    # Add new relationship columns for legacy databases
    ensure_column("task", 'assignee_id TEXT REFERENCES person(id) ON DELETE SET NULL')
    ensure_column("subtask", 'assignee_id TEXT REFERENCES person(id) ON DELETE SET NULL')
    ensure_column("activity", 'author_id TEXT REFERENCES person(id) ON DELETE SET NULL')
    ensure_column("activity", "task_context TEXT")
    ensure_column("project", "stakeholders JSON")
    ensure_column("project", "executiveUpdate TEXT")
    ensure_column("project", "startDate TEXT")
    ensure_column("project", "targetDate TEXT")
    ensure_column("project", "lastUpdate TEXT")
    ensure_column("project", "priority TEXT")
    ensure_column("project", "progress INTEGER")
    ensure_column("project", "status TEXT")
    ensure_column("project", "description TEXT")
    ensure_column("project", "initiative_id TEXT REFERENCES initiative(id) ON DELETE SET NULL")


def get_session():
    with Session(engine) as session:
        yield session


def extract_bearer_token(auth_header: str | None) -> str | None:
    if not auth_header:
        return None
    if not auth_header.lower().startswith("bearer "):
        return None
    return auth_header.split(None, 1)[1].strip() or None


def ensure_admin(request: Request) -> None:
    expected_token = os.getenv(ADMIN_TOKEN_ENV)
    if not expected_token:
        if current_environment() in PROTECTED_ENVIRONMENTS:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="Admin token not configured for this environment",
            )
        logger.warning(
            "Admin token not configured; allowing admin endpoints in non-protected environment"
        )
        return

    provided_token = (
        request.headers.get("x-admin-token")
        or extract_bearer_token(request.headers.get("authorization"))
    )
    if not provided_token or provided_token != expected_token:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Invalid admin token",
        )


def serialize_subtask(subtask: Subtask) -> dict:
    return {
        "id": subtask.id,
        "title": subtask.title,
        "status": subtask.status,
        "dueDate": subtask.dueDate,
        "completedDate": subtask.completedDate,
        "assigneeId": subtask.assignee_id,
        "assignee": serialize_person(subtask.assignee),
    }


def serialize_task(task: Task) -> dict:
    return {
        "id": task.id,
        "title": task.title,
        "status": task.status,
        "dueDate": task.dueDate,
        "completedDate": task.completedDate,
        "assigneeId": task.assignee_id,
        "assignee": serialize_person(task.assignee),
        "subtasks": [serialize_subtask(st) for st in task.subtasks],
    }


def serialize_activity(activity: Activity, person_index: PersonIndex | None = None) -> dict:
    import json
    task_context = None
    if activity.task_context:
        try:
            task_context = json.loads(activity.task_context)
        except (json.JSONDecodeError, TypeError):
            task_context = None

    resolved_person: Person | None = None
    if activity.author_person:
        resolved_person = activity.author_person
    elif person_index:
        resolved_person = person_index.resolve(
            person_id=activity.author_id,
            name=activity.author,
        )

    return {
        "id": activity.id,
        "date": activity.date,
        "note": activity.note,
        "taskContext": task_context,
        "author": (resolved_person.name if resolved_person else None) or activity.author,
        "authorId": resolved_person.id if resolved_person else activity.author_id,
        "authorPerson": serialize_person(resolved_person),
    }


def serialize_stakeholder(raw_stakeholder: dict, person_index: PersonIndex | None = None) -> dict:
    stakeholder = normalize_stakeholders([raw_stakeholder])[0]
    person: Person | None = None
    if person_index:
        person = person_index.resolve(
            name=stakeholder.get("name"),
            email=stakeholder.get("email"),
            person_id=stakeholder.get("id"),
        )

    if person:
        stakeholder["id"] = person.id
        stakeholder["name"] = person.name
        stakeholder["team"] = stakeholder.get("team") or person.team or ""
        stakeholder["email"] = person.email or stakeholder.get("email")

    return stakeholder


def normalize_project_activity(project: Project) -> Project:
    if project.recentActivity is None:
        project.recentActivity = []

    for activity in project.recentActivity:
        if not activity.author and activity.author_person:
            activity.author = activity.author_person.name

    project.recentActivity.sort(key=lambda a: a.date or "", reverse=True)

    if project.recentActivity:
        project.lastUpdate = project.recentActivity[0].note

    return project


def migrate_people_links(session: Session) -> None:
    """
    Convert legacy stakeholder/author data into normalized person relationships.
    """
    projects = session.exec(
        select(Project).options(
            selectinload(Project.stakeholders),
            selectinload(Project.plan).selectinload(Task.subtasks),
            selectinload(Project.plan).selectinload(Task.assignee),
            selectinload(Project.recentActivity).selectinload(Activity.author_person),
        )
    ).all()
    updated = False

    for project in projects:
        legacy_stakeholders = normalize_stakeholders(project.stakeholders_legacy)
        if legacy_stakeholders:
            for stakeholder in legacy_stakeholders:
                person = resolve_person_reference(session, stakeholder)
                if person and person not in project.stakeholders:
                    project.stakeholders.append(person)
                    updated = True
            project.stakeholders_legacy = []

        for activity in project.recentActivity or []:
            if activity.author_id is None and activity.author:
                person = resolve_person_reference(session, activity.author)
                if person:
                    activity.author_id = person.id
                    activity.author = person.name
                    updated = True

        for task in project.plan or []:
            if task.assignee_id and task.assignee is None:
                person = session.get(Person, task.assignee_id)
                if person:
                    task.assignee = person
                else:
                    task.assignee_id = None
                updated = True
            for subtask in task.subtasks or []:
                if subtask.assignee_id and subtask.assignee is None:
                    person = session.get(Person, subtask.assignee_id)
                    if person:
                        subtask.assignee = person
                    else:
                        subtask.assignee_id = None
                    updated = True

    if updated:
        session.commit()


def serialize_project(project: Project, person_index: PersonIndex | None = None) -> dict:
    normalize_project_activity(project)

    result = {
        "id": project.id,
        "name": project.name,
        "status": project.status,
        "priority": project.priority,
        "progress": project.progress,
        "lastUpdate": project.lastUpdate,
        "description": project.description,
        "executiveUpdate": project.executiveUpdate,
        "startDate": project.startDate,
        "targetDate": project.targetDate,
        "stakeholders": [serialize_person(person) for person in project.stakeholders],
        "plan": [serialize_task(task) for task in project.plan],
        "recentActivity": [serialize_activity(activity, person_index) for activity in project.recentActivity],
        "initiativeId": project.initiative_id,
    }
    # Include initiative name if available
    if project.initiative:
        result["initiative"] = {
            "id": project.initiative.id,
            "name": project.initiative.name,
        }
    else:
        result["initiative"] = None
    return result


def resolve_assignee_id(session: Session, assignee_payload: Optional[AssigneePayload]) -> Optional[str]:
    """Resolve an assignee payload to a person ID, looking up by id or name"""
    if assignee_payload is None:
        return None

    # If ID is provided, verify it exists
    if assignee_payload.id:
        person = session.exec(select(Person).where(Person.id == assignee_payload.id)).first()
        if person:
            return person.id

    # If name is provided, look up by name
    if assignee_payload.name:
        person = session.exec(select(Person).where(
            func.lower(Person.name) == assignee_payload.name.lower()
        )).first()
        if person:
            return person.id
        # Create new person if not found
        new_person = Person(
            id=generate_id("person"),
            name=assignee_payload.name,
            team=assignee_payload.team or "Contributor"
        )
        session.add(new_person)
        session.commit()
        session.refresh(new_person)
        return new_person.id

    return None

def serialize_project_with_people(session: Session, project: Project) -> dict:
    person_index = build_person_index(session)
    return serialize_project(project, person_index)


def serialize_initiative(initiative: Initiative, person_index: PersonIndex | None = None, include_projects: bool = True) -> dict:
    """Serialize an initiative to a dictionary for API response."""
    # Collect aggregated stakeholders from all projects (unique by id)
    stakeholder_ids = set()
    aggregated_stakeholders = []
    for project in initiative.projects:
        for stakeholder in project.stakeholders:
            if stakeholder.id not in stakeholder_ids:
                stakeholder_ids.add(stakeholder.id)
                aggregated_stakeholders.append(serialize_person(stakeholder))

    result = {
        "id": initiative.id,
        "name": initiative.name,
        "description": initiative.description,
        "status": initiative.status,
        "priority": initiative.priority,
        "startDate": initiative.startDate,
        "targetDate": initiative.targetDate,
        "owners": [serialize_person(owner) for owner in initiative.owners],
        "stakeholders": aggregated_stakeholders,
    }

    if include_projects:
        result["projects"] = [
            {
                "id": project.id,
                "name": project.name,
                "status": project.status,
                "priority": project.priority,
                "progress": project.progress,
            }
            for project in initiative.projects
        ]

    return result


def serialize_initiative_with_full_projects(initiative: Initiative, person_index: PersonIndex | None = None) -> dict:
    """Serialize an initiative with full project details."""
    result = serialize_initiative(initiative, person_index, include_projects=False)
    result["projects"] = [serialize_project(project, person_index) for project in initiative.projects]
    return result


def apply_task_payload(task: Task, payload: TaskPayload, session: Session | None = None) -> Task:
    task.title = payload.title
    task.status = payload.status
    task.dueDate = payload.dueDate
    task.completedDate = payload.completedDate

    # --- Task assignee ---
    if session is not None:
        fields_set = getattr(payload, "model_fields_set", getattr(payload, "__fields_set__", set()))
        # If payload explicitly includes the assignee field and it is None, clear it.
        if "assignee" in fields_set and payload.assignee is None:
            task.assignee = None
            task.assignee_id = None
        else:
            ref = None
            if "assignee" in fields_set and payload.assignee:
                ref = payload.assignee
            elif hasattr(payload, "assignee_id") and payload.assignee_id:
                ref = payload.assignee_id

            if ref is not None:
                assignee = resolve_person_reference(session, ref)
                task.assignee = assignee
                task.assignee_id = assignee.id if assignee else None

    # --- Subtasks ---
    if payload.subtasks is not None:
        if task.subtasks is None:
            task.subtasks = []
        else:
            task.subtasks.clear()

        for subtask_payload in payload.subtasks:
            subtask = Subtask(
                id=subtask_payload.id or generate_id("subtask"),
                title=subtask_payload.title,
                status=subtask_payload.status,
                dueDate=subtask_payload.dueDate,
                completedDate=subtask_payload.completedDate,
                assignee_id=subtask_payload.assignee_id,
            )

            if session is not None:
                # Same explicit-clear behavior for subtasks if 'assignee' exists
                subtask_fields_set = getattr(subtask_payload, "model_fields_set", getattr(subtask_payload, "__fields_set__", set()))
                if "assignee" in subtask_fields_set and subtask_payload.assignee is None:
                    subtask.assignee = None
                    subtask.assignee_id = None
                else:
                    ref = None
                    if "assignee" in subtask_fields_set and subtask_payload.assignee:
                        ref = subtask_payload.assignee
                    elif hasattr(subtask_payload, "assignee_id") and subtask_payload.assignee_id:
                        ref = subtask_payload.assignee_id

                    if ref is not None:
                        assignee = resolve_person_reference(session, ref)
                        subtask.assignee = assignee
                        subtask.assignee_id = assignee.id if assignee else None

            task.subtasks.append(subtask)

    return task


def normalize_project_stakeholders(session: Session, stakeholders: Optional[List[Stakeholder | dict]]) -> list[dict]:
    normalized: list[dict] = []
    for stakeholder in normalize_stakeholders(stakeholders):
        if not stakeholder.get("name"):
            continue

        person = upsert_person_from_details(
            session,
            name=stakeholder.get("name", ""),
            team=stakeholder.get("team", ""),
            email=stakeholder.get("email"),
            person_id=stakeholder.get("id"),
        )

        normalized.append(
            {
                "id": person.id,
                "name": person.name,
                "team": stakeholder.get("team") or person.team or "",
                "email": person.email,
            }
        )

    return normalized


def upsert_project(session: Session, payload: ProjectPayload) -> Project:
    normalized_name = (payload.name or "").strip()

    statement = (
        select(Project)
        .where(Project.id == payload.id)
        .options(
            selectinload(Project.stakeholders),
            selectinload(Project.plan).selectinload(Task.subtasks),
            selectinload(Project.plan).selectinload(Task.assignee),
            selectinload(Project.recentActivity).selectinload(Activity.author_person),
        )
    )
    project = session.exec(statement).first() if payload.id else None
    if project is None:
        project = Project(id=payload.id or generate_id("project"))

    # Check for duplicate project name (case-insensitive)
    existing_project = session.exec(
        select(Project).where(
            func.lower(Project.name) == func.lower(normalized_name),
            Project.id != (project.id if project else "")
        )
    ).first()
    if existing_project:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"A project with the name '{normalized_name}' already exists. Please choose a different name."
        )

    project.name = normalized_name
    project.status = payload.status
    project.priority = payload.priority
    project.progress = payload.progress
    project.lastUpdate = payload.lastUpdate
    project.description = payload.description
    project.executiveUpdate = payload.executiveUpdate
    project.startDate = payload.startDate
    project.targetDate = payload.targetDate
    project.stakeholders_legacy = []

    if project.stakeholders is None:
        project.stakeholders = []
    else:
        project.stakeholders.clear()
    seen_stakeholders: set[str] = set()
    for stakeholder_payload in payload.stakeholders:
        person = resolve_person_reference(session, stakeholder_payload)
        if person and person.id not in seen_stakeholders:
            project.stakeholders.append(person)
            seen_stakeholders.add(person.id)

    if project.plan is None:
        project.plan = []
    else:
        project.plan.clear()
    for task_payload in payload.plan:
        task = Task(
            id=task_payload.id or generate_id("task"),
            title=task_payload.title,
            status=task_payload.status,
            dueDate=task_payload.dueDate,
            completedDate=task_payload.completedDate,
            assignee_id=task_payload.assignee_id,
        )
        apply_task_payload(task, task_payload, session)
        project.plan.append(task)

    if project.recentActivity is None:
        project.recentActivity = []
    else:
        project.recentActivity.clear()
    activity_payloads = sorted(
        payload.recentActivity,
        key=lambda activity: activity.date or "",
    )

    import json as json_module
    for activity_payload in activity_payloads:
        # Serialize taskContext to JSON string if present
        task_context_str = None
        if activity_payload.taskContext is not None:
            task_context_str = json_module.dumps({
                "taskId": activity_payload.taskContext.taskId,
                "subtaskId": activity_payload.taskContext.subtaskId,
                "taskTitle": activity_payload.taskContext.taskTitle,
                "subtaskTitle": activity_payload.taskContext.subtaskTitle,
            })
        author_person = resolve_person_reference(session, activity_payload.author_id or activity_payload.author)
        author_name = (author_person.name if author_person else None) or activity_payload.author
        activity = Activity(
            id=activity_payload.id or generate_id("activity"),
            date=activity_payload.date,
            note=activity_payload.note,
            task_context=task_context_str,
            author=author_name,
            author_id=author_person.id if author_person else activity_payload.author_id,
        )
        project.recentActivity.append(activity)

    normalize_project_activity(project)

    session.add(project)
    session.commit()
    # Reload project with all relationships properly loaded
    return load_project(session, project.id)


def add_data_change_activity(
    session: Session,
    project_id: str,
    request: Request | None,
    note: str,
    author: str | None = None
) -> Activity:
    """Add an activity entry for a data change to the project's activity feed"""
    author_name, author_id = resolve_activity_author(session, request, author)
    activity = Activity(
        id=generate_id("activity"),
        date=datetime.utcnow().isoformat(),
        note=note,
        author=author_name,
        author_id=author_id,
        project_id=project_id,
    )
    session.add(activity)
    session.commit()

    # Normalize project activity to update lastUpdate
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if project:
        normalize_project_activity(project)
        session.add(project)
        session.commit()

    return activity
def run_people_backfill(session: Session) -> None:
    migration_key = "people-backfill-v1"
    if session.get(MigrationState, migration_key):
        return

    projects = session.exec(
        select(Project)
        .options(
            selectinload(Project.plan).selectinload(Task.subtasks),
            selectinload(Project.recentActivity),
        )
    ).all()

    for project in projects:
        updated = False

        legacy_stakeholders = [
            stakeholder for stakeholder in project.stakeholders_legacy or [] if not isinstance(stakeholder, Person)
        ]
        normalized_stakeholders = normalize_project_stakeholders(session, legacy_stakeholders)
        if normalized_stakeholders:
            existing_person_ids = {person.id for person in project.stakeholders if person.id}
            for stakeholder in normalized_stakeholders:
                person_id = stakeholder.get("id")
                if not person_id or person_id in existing_person_ids:
                    continue

                person = session.get(Person, person_id)
                if person is None:
                    continue

                project.stakeholders.append(person)
                existing_person_ids.add(person_id)
                updated = True

            project.stakeholders_legacy = []
            updated = True

        for activity in project.recentActivity or []:
            if not activity.author:
                continue
            person = upsert_person_from_details(session, name=activity.author)
            if person and activity.author != person.name:
                activity.author = person.name
                updated = True

        if updated:
            session.add(project)

    session.add(MigrationState(key=migration_key))
    session.commit()


@app.post("/admin/people-backfill", dependencies=[Depends(ensure_admin)])
def trigger_people_backfill(session: Session = Depends(get_session)) -> dict:
    run_people_backfill(session)
    return {"status": "ok", "message": "People backfill completed or already applied."}


@app.on_event("startup")
def on_startup() -> None:
    create_db_and_tables()

    # Run database migrations
    with Session(engine) as session:
        migrate_add_unique_constraints(session)
        migrate_remove_email_credentials(session)

    if not is_dev_seeding_enabled():
        return

    with Session(engine) as session:
        migrate_people_links(session)
        existing = session.exec(select(Project)).first()
        if existing:
            return

        default_projects = [
            ProjectPayload(
                name="Website Redesign",
                status="active",
                priority="high",
                progress=45,
                lastUpdate="Completed homepage mockups and testing",
                description="Overhaul of company site for better UX.",
                executiveUpdate="Overhaul of company site for better UX.",
                startDate="2025-10-15",
                targetDate="2025-12-20",
                stakeholders=[PersonReference(name="Sarah Chen", team="Design"), PersonReference(name="Marcus Rodriguez", team="Development")],
                plan=[
                    TaskPayload(
                        title="Discovery & Research",
                        status="completed",
                        dueDate="2025-11-01",
                        completedDate="2025-11-01",
                        subtasks=[
                            SubtaskPayload(title="Competitive analysis", status="completed", completedDate="2025-10-22", dueDate="2025-10-22"),
                            SubtaskPayload(title="User interviews", status="completed", completedDate="2025-10-28", dueDate="2025-10-28"),
                        ],
                    ),
                    TaskPayload(
                        title="Design Phase",
                        status="in-progress",
                        dueDate="2025-12-05",
                        subtasks=[
                            SubtaskPayload(title="Homepage mockups", status="completed", completedDate="2025-11-28", dueDate="2025-11-28"),
                            SubtaskPayload(title="Product page designs", status="in-progress", dueDate="2025-12-03"),
                        ],
                    ),
                ],
                recentActivity=[
                    ActivityPayload(date="2025-11-29T14:30:00", note="Positive feedback on homepage direction", author="Alex Morgan"),
                    ActivityPayload(date="2025-11-28T16:15:00", note="Completed homepage mockups", author="Alex Morgan"),
                ],
            ),
            ProjectPayload(
                name="Q4 Marketing Campaign",
                status="active",
                priority="medium",
                progress=30,
                lastUpdate="Draft content calendar completed",
                description="Multi-channel campaign for Q4.",
                executiveUpdate="Multi-channel campaign for Q4.",
                startDate="2025-11-01",
                targetDate="2025-12-31",
                stakeholders=[PersonReference(name="Jennifer Liu", team="Marketing"), PersonReference(name="Alex Thompson", team="Creative")],
                plan=[
                    TaskPayload(
                        title="Campaign Strategy",
                        status="completed",
                        dueDate="2025-11-15",
                        completedDate="2025-11-15",
                        subtasks=[
                            SubtaskPayload(title="Define target audience", status="completed", completedDate="2025-11-05", dueDate="2025-11-05"),
                            SubtaskPayload(title="Set campaign goals", status="completed", completedDate="2025-11-10", dueDate="2025-11-10"),
                        ],
                    ),
                ],
                recentActivity=[
                    ActivityPayload(date="2025-11-27T16:30:00", note="Met with marketing to discuss timeline", author="Alex Morgan"),
                ],
            ),
        ]

        for project_payload in default_projects:
            upsert_project(session, project_payload)


def load_project(session: Session, project_id: str) -> Project:
    statement = (
        select(Project)
        .where(Project.id == project_id)
        .options(
            selectinload(Project.plan).selectinload(Task.subtasks),
            selectinload(Project.plan).selectinload(Task.assignee),
            selectinload(Project.plan).selectinload(Task.subtasks).selectinload(Subtask.assignee),
            selectinload(Project.recentActivity),
            selectinload(Project.recentActivity).selectinload(Activity.author_person),
            selectinload(Project.stakeholders),
            selectinload(Project.initiative),
        )
    )
    project = session.exec(statement).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
    return project


@app.get("/people")
def list_people(session: Session = Depends(get_session)):
    statement = select(Person)
    people = session.exec(statement).all()

    unique_people: dict[str, Person] = {}
    seen_people: set[str] = set()
    for person in people:
        email_key = person.email.lower() if person.email else None
        name_key = person.name.lower() if person.name else None

        existing = None
        if email_key and email_key in unique_people:
            existing = unique_people[email_key]
        elif name_key and name_key in unique_people:
            existing = unique_people[name_key]

        if existing is None:
            if name_key:
                unique_people[name_key] = person
            if email_key:
                unique_people[email_key] = person
            continue

        # Collapse legacy duplicates by preferring the first encountered record
        existing.team = existing.team or person.team
        existing.email = existing.email or person.email
        session.delete(person)

    # Ensure each person appears only once even though we index by multiple keys
    deduped_people: list[Person] = []
    for person in unique_people.values():
        if person.id in seen_people:
            continue
        seen_people.add(person.id)
        deduped_people.append(person)

    session.commit()
    return [serialize_person(person) for person in deduped_people]


@app.post("/people", status_code=status.HTTP_201_CREATED)
def create_person(payload: PersonPayload, request: Request, session: Session = Depends(get_session)):
    person = upsert_person_from_payload(session, payload)
    log_action(session, "create_person", "person", person.id, {"name": person.name, "team": person.team}, request)
    return serialize_person(person)


@app.get("/people/{person_id}")
def get_person(person_id: str, session: Session = Depends(get_session)):
    person = session.exec(select(Person).where(Person.id == person_id)).first()
    if not person:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Person not found")
    return serialize_person(person)


@app.put("/people/{person_id}")
def update_person(person_id: str, payload: PersonPayload, request: Request, session: Session = Depends(get_session)):
    person = session.exec(select(Person).where(Person.id == person_id)).first()
    if not person:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Person not found")

    normalized_name = payload.name.strip()

    conflict = get_person_by_name(session, normalized_name)
    if conflict and conflict.id != person.id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="A person with that name already exists",
        )

    old_values = {"name": person.name, "team": person.team, "email": person.email}
    person.name = normalized_name
    person.team = payload.team
    person.email = payload.email
    session.add(person)
    session.commit()
    session.refresh(person)
    log_action(session, "update_person", "person", person_id, {"old": old_values, "new": {"name": person.name, "team": person.team, "email": person.email}}, request)
    return serialize_person(person)


@app.delete("/people/{person_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_person(person_id: str, request: Request, session: Session = Depends(get_session)):
    person = session.exec(select(Person).where(Person.id == person_id)).first()
    if not person:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Person not found")
    deleted_data = {"name": person.name, "team": person.team}
    session.delete(person)
    session.commit()
    log_action(session, "delete_person", "person", person_id, deleted_data, request)
    return None


@app.get("/settings/email", response_model=EmailSettingsResponse)
def read_email_settings(session: Session = Depends(get_session)):
    settings = get_email_settings(session)
    return serialize_email_settings(settings)


@app.put("/settings/email", response_model=EmailSettingsResponse)
def update_email_settings(payload: EmailSettingsPayload, session: Session = Depends(get_session)):
    settings = get_email_settings(session)
    settings.smtp_server = payload.smtpServer
    settings.smtp_port = payload.smtpPort
    settings.from_address = payload.fromAddress or None
    settings.use_tls = payload.useTLS

    session.add(settings)
    session.commit()
    session.refresh(settings)
    return serialize_email_settings(settings)


@app.post("/actions/email", status_code=status.HTTP_202_ACCEPTED)
def send_email_action(payload: EmailSendPayload, request: Request, session: Session = Depends(get_session)):
    """
    Send an email and verify server acceptance.

    SMTP settings are now passed inline with each request (stored in browser localStorage).
    Returns 200 with sent_to list on success.
    Returns 400 for configuration/validation errors.
    Returns 502 for SMTP server errors.
    """
    recipients = normalize_recipients(payload.recipients)
    cc = normalize_recipients(payload.cc or [])
    bcc = normalize_recipients(payload.bcc or [])
    recipient_count = len(recipients) + len(cc) + len(bcc)

    settings = get_email_settings(session)

    smtp_server = payload.smtp_server or settings.smtp_server
    smtp_port = payload.smtp_port or settings.smtp_port
    from_address = payload.from_address or settings.from_address
    use_tls = payload.use_tls if payload.use_tls is not None else settings.use_tls

    try:
        result = dispatch_email(
            smtp_server=smtp_server or "",
            smtp_port=smtp_port or 587,
            from_address=from_address or "",
            recipients=recipients,
            cc=cc,
            bcc=bcc,
            subject=payload.subject,
            body=payload.body,
            use_tls=use_tls if use_tls is not None else True
        )
    except ValueError as exc:
        log_action(session, "send_email_failed", "email", None, {"error": str(exc), "recipient_count": recipient_count}, request)
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc))
    except smtplib.SMTPException as exc:
        log_action(session, "send_email_failed", "email", None, {"error": str(exc), "recipient_count": recipient_count}, request)
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail=f"SMTP error: {str(exc)}")

    log_action(session, "send_email", "email", None, {"recipient_count": len(result["sent_to"]), "subject_preview": payload.subject[:50] if payload.subject else None}, request)

    # Return detailed result including which recipients were successful
    return {
        "status": "sent",
        "sent_to": result["sent_to"],
        "refused": result["refused"],
        "message": f"Email delivered to {len(result['sent_to'])} recipient(s)"
    }


# =============================================================================
# Initiative Endpoints
# =============================================================================

def load_initiative(session: Session, initiative_id: str) -> Initiative:
    """Load an initiative by ID with all relationships."""
    statement = select(Initiative).where(Initiative.id == initiative_id).options(
        selectinload(Initiative.projects).selectinload(Project.stakeholders),
        selectinload(Initiative.owners),
    )
    initiative = session.exec(statement).first()
    if not initiative:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Initiative not found")
    return initiative


def upsert_initiative(session: Session, payload: InitiativePayload) -> Initiative:
    """Create or update an initiative from a payload."""
    initiative_id = payload.id or generate_id("initiative")

    # Check if initiative exists
    existing = session.exec(select(Initiative).where(Initiative.id == initiative_id)).first()

    if existing:
        initiative = existing
        initiative.name = payload.name
        initiative.description = payload.description
        initiative.status = payload.status
        initiative.priority = payload.priority
        initiative.startDate = payload.startDate
        initiative.targetDate = payload.targetDate
    else:
        initiative = Initiative(
            id=initiative_id,
            name=payload.name,
            description=payload.description,
            status=payload.status,
            priority=payload.priority,
            startDate=payload.startDate,
            targetDate=payload.targetDate,
        )
        session.add(initiative)

    # Handle owners
    if payload.owners:
        # Clear existing owners and add new ones
        initiative.owners.clear()
        for owner_ref in payload.owners:
            person = resolve_person_reference(session, owner_ref)
            if person and person not in initiative.owners:
                initiative.owners.append(person)

    session.commit()
    session.refresh(initiative)
    return initiative


@app.get("/initiatives")
def list_initiatives(session: Session = Depends(get_session)):
    """List all initiatives with their projects and owners."""
    person_index = build_person_index(session)
    statement = select(Initiative).options(
        selectinload(Initiative.projects).selectinload(Project.stakeholders),
        selectinload(Initiative.owners),
    )
    initiatives = session.exec(statement).all()
    return [serialize_initiative(initiative, person_index) for initiative in initiatives]


@app.post("/initiatives", status_code=status.HTTP_201_CREATED)
def create_initiative(payload: InitiativePayload, request: Request, session: Session = Depends(get_session)):
    """Create a new initiative."""
    initiative = upsert_initiative(session, payload)
    log_action(session, "create_initiative", "initiative", initiative.id, {"name": initiative.name}, request)
    # Reload with relationships
    initiative = load_initiative(session, initiative.id)
    return serialize_initiative(initiative)


@app.get("/initiatives/{initiative_id}")
def get_initiative(initiative_id: str, session: Session = Depends(get_session)):
    """Get a single initiative with full details."""
    person_index = build_person_index(session)
    initiative = load_initiative(session, initiative_id)
    return serialize_initiative_with_full_projects(initiative, person_index)


@app.put("/initiatives/{initiative_id}")
def update_initiative(initiative_id: str, payload: InitiativePayload, request: Request, session: Session = Depends(get_session)):
    """Update an initiative."""
    if payload.id and payload.id != initiative_id:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Initiative ID mismatch")

    payload.id = initiative_id
    initiative = upsert_initiative(session, payload)
    log_action(session, "update_initiative", "initiative", initiative_id, {"name": initiative.name}, request)
    # Reload with relationships
    initiative = load_initiative(session, initiative.id)
    return serialize_initiative(initiative)


@app.delete("/initiatives/{initiative_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_initiative(initiative_id: str, request: Request, session: Session = Depends(get_session)):
    """Delete an initiative. Projects become ungrouped (not deleted)."""
    initiative = load_initiative(session, initiative_id)
    deleted_data = {"name": initiative.name}

    # Projects will become ungrouped due to ON DELETE SET NULL
    session.delete(initiative)
    session.commit()
    log_action(session, "delete_initiative", "initiative", initiative_id, deleted_data, request)
    return None


@app.post("/initiatives/{initiative_id}/owners/{person_id}", status_code=status.HTTP_201_CREATED)
def add_owner_to_initiative(initiative_id: str, person_id: str, request: Request, session: Session = Depends(get_session)):
    """Add an owner to an initiative."""
    initiative = load_initiative(session, initiative_id)
    person = session.exec(select(Person).where(Person.id == person_id)).first()
    if not person:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Person not found")

    if person not in initiative.owners:
        initiative.owners.append(person)
        session.commit()
        log_action(session, "add_initiative_owner", "initiative", initiative_id, {"person_id": person_id, "person_name": person.name}, request)

    return serialize_initiative(initiative)


@app.delete("/initiatives/{initiative_id}/owners/{person_id}", status_code=status.HTTP_204_NO_CONTENT)
def remove_owner_from_initiative(initiative_id: str, person_id: str, request: Request, session: Session = Depends(get_session)):
    """Remove an owner from an initiative."""
    initiative = load_initiative(session, initiative_id)
    person = session.exec(select(Person).where(Person.id == person_id)).first()
    if not person:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Person not found")

    if person in initiative.owners:
        initiative.owners.remove(person)
        session.commit()
        log_action(session, "remove_initiative_owner", "initiative", initiative_id, {"person_id": person_id, "person_name": person.name}, request)

    return None


@app.post("/initiatives/{initiative_id}/projects/{project_id}", status_code=status.HTTP_201_CREATED)
def add_project_to_initiative(initiative_id: str, project_id: str, request: Request, session: Session = Depends(get_session)):
    """Add a project to an initiative."""
    initiative = load_initiative(session, initiative_id)
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")

    project.initiative_id = initiative_id
    session.commit()
    log_action(session, "add_project_to_initiative", "initiative", initiative_id, {"project_id": project_id, "project_name": project.name}, request)

    # Reload initiative
    initiative = load_initiative(session, initiative_id)
    return serialize_initiative(initiative)


@app.delete("/initiatives/{initiative_id}/projects/{project_id}", status_code=status.HTTP_204_NO_CONTENT)
def remove_project_from_initiative(initiative_id: str, project_id: str, request: Request, session: Session = Depends(get_session)):
    """Remove a project from an initiative (makes it ungrouped)."""
    initiative = load_initiative(session, initiative_id)
    project = session.exec(select(Project).where(Project.id == project_id)).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")

    if project.initiative_id == initiative_id:
        project.initiative_id = None
        session.commit()
        log_action(session, "remove_project_from_initiative", "initiative", initiative_id, {"project_id": project_id, "project_name": project.name}, request)

    return None


@app.get("/projects")
def list_projects(session: Session = Depends(get_session)):
    person_index = build_person_index(session)
    statement = select(Project).options(
        selectinload(Project.plan).selectinload(Task.subtasks),
        selectinload(Project.plan).selectinload(Task.assignee),
        selectinload(Project.plan).selectinload(Task.subtasks).selectinload(Subtask.assignee),
        selectinload(Project.recentActivity),
        selectinload(Project.recentActivity).selectinload(Activity.author_person),
        selectinload(Project.stakeholders),
        selectinload(Project.initiative),
    )
    projects = session.exec(statement).all()
    return [serialize_project(project, person_index) for project in projects]


@app.post("/projects", status_code=status.HTTP_201_CREATED)
def create_project(payload: ProjectPayload, request: Request, session: Session = Depends(get_session)):
    project = upsert_project(session, payload)
    log_action(session, "create_project", "project", project.id, {"name": project.name, "status": project.status}, request)
    return serialize_project_with_people(session, project)


@app.get("/projects/{project_id}")
def get_project(project_id: str, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    return serialize_project_with_people(session, project)


@app.put("/projects/{project_id}")
def update_project(project_id: str, payload: ProjectPayload, request: Request, session: Session = Depends(get_session)):
    if payload.id and payload.id != project_id:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Project ID mismatch")

    # Get existing project to track changes
    existing = session.exec(select(Project).where(Project.id == project_id)).first()
    old_values = {}
    if existing:
        old_values = {
            "name": existing.name,
            "status": existing.status,
            "priority": existing.priority,
            "progress": existing.progress,
            "description": existing.description,
            "executiveUpdate": existing.executiveUpdate,
            "startDate": existing.startDate,
            "targetDate": existing.targetDate,
        }

    payload.id = project_id
    project = upsert_project(session, payload)

    # Build activity description with specific changes
    changes = []
    if old_values:
        if old_values["name"] != project.name:
            changes.append(f"name to: {project.name}")
        if old_values["status"] != project.status:
            changes.append(f"status to: {project.status}")
        if old_values["priority"] != project.priority:
            changes.append(f"priority to: {project.priority}")
        if old_values["progress"] != project.progress:
            changes.append(f"progress to: {project.progress}%")
        if old_values["description"] != project.description:
            changes.append(f"description to: {project.description[:100]}{'...' if len(project.description or '') > 100 else ''}")
        if old_values["executiveUpdate"] != project.executiveUpdate:
            changes.append(f"executive update to: {(project.executiveUpdate or '')[:100]}{'...' if len(project.executiveUpdate or '') > 100 else ''}")
        if old_values["startDate"] != project.startDate:
            changes.append(f"start date to: {project.startDate or 'none'}")
        if old_values["targetDate"] != project.targetDate:
            changes.append(f"target date to: {project.targetDate or 'none'}")

        if changes:
            add_data_change_activity(
                session, project_id, request,
                f"Updated project: {', '.join(changes)}"
            )

    log_action(session, "update_project", "project", project_id, {"name": project.name, "status": project.status}, request)
    return serialize_project_with_people(session, project)


@app.delete("/projects/{project_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_project(project_id: str, request: Request, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    deleted_data = {"name": project.name}
    session.delete(project)
    session.commit()
    log_action(session, "delete_project", "project", project_id, deleted_data, request)
    return None


@app.post("/projects/{project_id}/tasks")
def create_task(project_id: str, payload: TaskPayload, request: Request, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    task = Task(
        id=payload.id or generate_id("task"),
        title=payload.title,
        status=payload.status,
        dueDate=payload.dueDate,
        completedDate=payload.completedDate,
        project_id=project.id,
        assignee_id=payload.assignee_id,
    )
    apply_task_payload(task, payload, session)
    session.add(task)
    session.commit()
    session.refresh(task)

    log_action(session, "create_task", "task", task.id, {"project_id": project_id, "title": task.title}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.put("/projects/{project_id}/tasks/{task_id}")
def update_task(project_id: str, task_id: str, payload: TaskPayload, request: Request, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    task = session.exec(select(Task).where(Task.id == task_id, Task.project_id == project.id)).first()
    if not task:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Task not found")

    # Track changes for activity feed
    changes = []
    old_status = task.status
    old_title = task.title
    old_due_date = task.dueDate
    old_assignee_id = task.assignee_id
    apply_task_payload(task, payload, session)
    session.add(task)
    session.commit()

    # Build activity description with specific changes
    if old_title != task.title:
        changes.append(f"title to: {task.title}")
    if old_status != task.status:
        changes.append(f"status from {old_status} to {task.status}")
    if old_due_date != task.dueDate:
        changes.append(f"due date to: {task.dueDate or 'none'}")
    if old_assignee_id != task.assignee_id:
        if task.assignee_id:
            new_assignee = session.exec(select(Person).where(Person.id == task.assignee_id)).first()
            changes.append(f"assigned to {new_assignee.name if new_assignee else 'unknown'}")
        else:
            changes.append("unassigned")

    if changes and not (len(changes) == 1 and old_title != task.title):
        add_data_change_activity(
            session, project_id, request,
            f"Updated task '{task.title}': {', '.join(changes)}"
        )

    log_action(session, "update_task", "task", task_id, {"project_id": project_id, "title": task.title, "old_status": old_status, "new_status": task.status}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.delete("/projects/{project_id}/tasks/{task_id}", status_code=status.HTTP_204_NO_CONTENT)
def remove_task(project_id: str, task_id: str, request: Request, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    task = session.exec(select(Task).where(Task.id == task_id, Task.project_id == project.id)).first()
    if not task:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Task not found")
    deleted_data = {"project_id": project_id, "title": task.title}
    task_title = task.title
    session.delete(task)
    session.commit()

    # Add activity for task deletion
    add_data_change_activity(session, project_id, request, f"Deleted task: {task_title}")

    log_action(session, "delete_task", "task", task_id, deleted_data, request)
    return None


@app.post("/projects/{project_id}/tasks/{task_id}/subtasks")
def create_subtask(project_id: str, task_id: str, payload: SubtaskPayload, request: Request, session: Session = Depends(get_session)):
    load_project(session, project_id)
    task = session.exec(select(Task).where(Task.id == task_id, Task.project_id == project_id)).first()
    if not task:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Task not found")
    assignee = resolve_person_reference(session, payload.assignee or payload.assignee_id)
    subtask = Subtask(
        id=payload.id or generate_id("subtask"),
        title=payload.title,
        status=payload.status,
        dueDate=payload.dueDate,
        completedDate=payload.completedDate,
        task_id=task.id,
        assignee_id=assignee.id if assignee else payload.assignee_id,
    )
    if assignee:
        subtask.assignee = assignee
    session.add(subtask)
    session.commit()

    log_action(session, "create_subtask", "subtask", subtask.id, {"project_id": project_id, "task_id": task_id, "title": subtask.title}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.put("/projects/{project_id}/tasks/{task_id}/subtasks/{subtask_id}")
def update_subtask(project_id: str, task_id: str, subtask_id: str, payload: SubtaskPayload, request: Request, session: Session = Depends(get_session)):
    load_project(session, project_id)
    subtask = session.exec(select(Subtask).where(Subtask.id == subtask_id, Subtask.task_id == task_id)).first()
    if not subtask:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Subtask not found")

    # Get parent task for activity message
    task = session.exec(select(Task).where(Task.id == task_id)).first()
    task_title = task.title if task else "Unknown Task"

    # Track changes for activity feed
    changes = []
    old_status = subtask.status
    old_title = subtask.title
    old_due_date = subtask.dueDate
    old_assignee_id = subtask.assignee_id
    assignee = resolve_person_reference(session, payload.assignee or payload.assignee_id)
    subtask.title = payload.title
    subtask.status = payload.status
    subtask.dueDate = payload.dueDate
    subtask.completedDate = payload.completedDate
    # If caller explicitly provides assignee=None, treat that as "clear", regardless of assignee_id.
    fields_set = getattr(payload, "model_fields_set", getattr(payload, "__fields_set__", set()))
    if "assignee" in fields_set and payload.assignee is None:
        subtask.assignee = None
        subtask.assignee_id = None
    else:
        subtask.assignee = assignee
        subtask.assignee_id = assignee.id if assignee else payload.assignee_id
    
    session.add(subtask)
    session.commit()

    # Build activity description with specific changes
    if old_title != subtask.title:
        changes.append(f"title to: {subtask.title}")
    if old_status != subtask.status:
        changes.append(f"status from {old_status} to {subtask.status}")
    if old_due_date != subtask.dueDate:
        changes.append(f"due date to: {subtask.dueDate or 'none'}")
    if old_assignee_id != subtask.assignee_id:
        if subtask.assignee_id:
            new_assignee = session.exec(select(Person).where(Person.id == subtask.assignee_id)).first()
            changes.append(f"assigned to {new_assignee.name if new_assignee else 'unknown'}")
        else:
            changes.append("unassigned")

    if changes and not (len(changes) == 1 and old_title != subtask.title):
        add_data_change_activity(
            session, project_id, request,
            f"Updated subtask '{subtask.title}' in task '{task_title}': {', '.join(changes)}"
        )

    log_action(session, "update_subtask", "subtask", subtask_id, {"project_id": project_id, "task_id": task_id, "title": subtask.title, "old_status": old_status, "new_status": subtask.status}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.delete("/projects/{project_id}/tasks/{task_id}/subtasks/{subtask_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_subtask(project_id: str, task_id: str, subtask_id: str, request: Request, session: Session = Depends(get_session)):
    load_project(session, project_id)
    subtask = session.exec(select(Subtask).where(Subtask.id == subtask_id, Subtask.task_id == task_id)).first()
    if not subtask:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Subtask not found")

    # Get parent task for activity message
    task = session.exec(select(Task).where(Task.id == task_id)).first()
    task_title = task.title if task else "Unknown Task"

    deleted_data = {"project_id": project_id, "task_id": task_id, "title": subtask.title}
    subtask_title = subtask.title
    session.delete(subtask)
    session.commit()

    # Add activity for subtask deletion
    add_data_change_activity(session, project_id, request, f"Deleted subtask '{subtask_title}' from task '{task_title}'")

    log_action(session, "delete_subtask", "subtask", subtask_id, deleted_data, request)
    return None


@app.post("/projects/{project_id}/activities")
def create_activity(project_id: str, payload: ActivityPayload, request: Request, session: Session = Depends(get_session)):
    import json as json_module
    project = load_project(session, project_id)

    # Serialize taskContext to JSON string if present
    task_context_str = None
    if payload.taskContext is not None:
        task_context_str = json_module.dumps({
            "taskId": payload.taskContext.taskId,
            "subtaskId": payload.taskContext.subtaskId,
            "taskTitle": payload.taskContext.taskTitle,
            "subtaskTitle": payload.taskContext.subtaskTitle,
        })

    author_person = resolve_person_reference(session, payload.author_id or payload.author)
    author_name = (author_person.name if author_person else None) or payload.author or "Unknown"
    activity = Activity(
        id=payload.id or generate_id("activity"),
        date=payload.date,
        note=payload.note,
        author=author_name,
        author_id=author_person.id if author_person else payload.author_id,
        project_id=project.id,
        task_context=task_context_str,
    )
    session.add(activity)
    session.commit()
    project = load_project(session, project_id)
    normalize_project_activity(project)
    session.add(project)
    session.commit()
    log_action(session, "create_activity", "activity", activity.id, {"project_id": project_id, "author": activity.author, "note_preview": activity.note[:100] if activity.note else None}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.put("/projects/{project_id}/activities/{activity_id}")
def update_activity(project_id: str, activity_id: str, payload: ActivityPayload, request: Request, session: Session = Depends(get_session)):
    import json as json_module
    load_project(session, project_id)
    activity = session.exec(select(Activity).where(Activity.id == activity_id, Activity.project_id == project_id)).first()
    if not activity:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Activity not found")
    author_person = resolve_person_reference(session, payload.author_id or payload.author)
    activity.note = payload.note
    activity.date = payload.date
    # Update taskContext if provided
    fields_set = getattr(payload, "model_fields_set", None) or getattr(payload, "__fields_set__", set())
    if payload.taskContext is not None:
        activity.task_context = json_module.dumps({
            "taskId": payload.taskContext.taskId,
            "subtaskId": payload.taskContext.subtaskId,
            "taskTitle": payload.taskContext.taskTitle,
            "subtaskTitle": payload.taskContext.subtaskTitle,
        })
    elif "taskContext" in fields_set:
        activity.task_context = None
    activity.author = (author_person.name if author_person else None) or payload.author or "Unknown"
    activity.author_id = author_person.id if author_person else payload.author_id
    session.add(activity)
    session.commit()
    project = load_project(session, project_id)
    normalize_project_activity(project)
    session.add(project)
    session.commit()
    log_action(session, "update_activity", "activity", activity_id, {"project_id": project_id, "author": activity.author}, request)
    return serialize_project_with_people(session, load_project(session, project_id))


@app.delete("/projects/{project_id}/activities/{activity_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_activity(project_id: str, activity_id: str, request: Request, session: Session = Depends(get_session)):
    project = load_project(session, project_id)
    activity = session.exec(select(Activity).where(Activity.id == activity_id, Activity.project_id == project_id)).first()
    if not activity:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Activity not found")
    deleted_data = {"project_id": project_id, "author": activity.author}
    session.delete(activity)
    session.commit()
    project = load_project(session, project_id)
    normalize_project_activity(project)
    session.add(project)
    session.commit()
    log_action(session, "delete_activity", "activity", activity_id, deleted_data, request)
    return None


@app.get("/export")
def export_portfolio(project_id: Optional[str] = None, session: Session = Depends(get_session)):
    projects = []
    if project_id:
        projects.append(serialize_project_with_people(session, load_project(session, project_id)))
    else:
        projects = list_projects(session)

    people = list_people(session)

    def iter_payload():
        yield "{\n"
        yield f"  \"version\": 1,\n"
        yield f"  \"exportedAt\": \"{datetime.utcnow().isoformat()}\",\n"
        yield "  \"projects\": "
        import json
        yield json.dumps(projects, indent=2)
        yield ",\n"
        yield "  \"people\": "
        yield json.dumps(people, indent=2)
        yield "\n}"

    return StreamingResponse(iter_payload(), media_type="application/json")


@app.post("/import")
async def import_portfolio(
    request: Request,
    payload: ImportPayload | None = Body(None),
    file: UploadFile | None = File(None),
    mode: str = "replace",
    session: Session = Depends(get_session),
):
    resolved_mode = mode

    if payload is None:
        if file is not None:
            try:
                import json

                data = json.loads(file.file.read())
                if "mode" not in data:
                    data["mode"] = resolved_mode
                payload = ImportPayload(**data)
            except Exception as exc:  # pragma: no cover - defensive
                raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Invalid import file: {exc}")
        else:
            try:
                import json

                raw_body = await request.body()
                if raw_body:
                    data = json.loads(raw_body)
                    if "mode" not in data:
                        data["mode"] = resolved_mode
                    payload = ImportPayload(**data)
            except Exception as exc:  # pragma: no cover - defensive
                raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Invalid import file: {exc}")

    if payload is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No import payload provided")

    if "mode" not in payload.model_fields_set:
        payload = payload.model_copy(update={"mode": resolved_mode})

    if payload.mode not in {"replace", "merge"}:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid import mode")

    existing_projects = {project.id: project for project in session.exec(select(Project)).all()}
    existing_people = {person.id: person for person in session.exec(select(Person)).all()}

    if payload.mode == "replace":
        session.exec(delete(Subtask))
        session.exec(delete(Task))
        session.exec(delete(Activity))
        session.exec(delete(Project))
        session.exec(delete(Person))
        session.commit()
        existing_projects = {}
        existing_people = {}

    for project_payload in payload.projects:
        if payload.mode == "merge" and project_payload.id in existing_projects:
            session.delete(existing_projects[project_payload.id])
            session.commit()
        upsert_project(session, project_payload)

    for person_payload in payload.people:
        if payload.mode == "merge" and person_payload.id in existing_people:
            session.delete(existing_people[person_payload.id])
            session.commit()

        upsert_person_from_payload(session, person_payload)

    projects = list_projects(session)
    people = list_people(session)

    log_action(session, "import_portfolio", "portfolio", None, {"mode": payload.mode, "project_count": len(payload.projects), "people_count": len(payload.people)}, request)

    return {"projects": projects, "people": people}


def _resolve_provider(provider_override: ChatProvider | None) -> ChatProvider:
    if provider_override:
        return provider_override

    env_provider = os.getenv("LLM_PROVIDER", ChatProvider.OPENAI.value)
    try:
        return ChatProvider(env_provider.lower())
    except ValueError as exc:  # pragma: no cover - defensive
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Unsupported LLM provider configured: {env_provider}",
        ) from exc


def _serialize_message(message: ChatMessage) -> dict:
    """Serialize a ChatMessage to OpenAI API format, excluding None values."""
    result = {"role": message.role.value}

    if message.content is not None:
        result["content"] = message.content

    if message.tool_calls:
        result["tool_calls"] = [tc.model_dump() for tc in message.tool_calls]

    if message.tool_call_id:
        result["tool_call_id"] = message.tool_call_id

    if message.name:
        result["name"] = message.name

    return result


def _build_llm_request(payload: ChatRequest, stream: bool = False) -> tuple[str, dict, dict]:
    provider = _resolve_provider(payload.provider)

    request_body = {
        "model": payload.model,
        "messages": [_serialize_message(m) for m in payload.messages],
    }

    if stream:
        request_body["stream"] = True

    if payload.response_format is not None:
        request_body["response_format"] = payload.response_format

    if payload.tools:
        request_body["tools"] = [t.model_dump() for t in payload.tools]

    if payload.tool_choice:
        request_body["tool_choice"] = payload.tool_choice

    if provider is ChatProvider.OPENAI:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="OpenAI API key not configured on server",
            )

        base_url = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
        url = f"{base_url}/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
        return url, headers, request_body

    api_key = os.getenv("AZURE_OPENAI_API_KEY")
    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", "").rstrip("/")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
    api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")

    if not api_key or not endpoint or not deployment:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Azure OpenAI configuration is incomplete",
        )

    url = (
        f"{endpoint}/openai/deployments/{deployment}/chat/completions"
        f"?api-version={api_version}"
    )
    headers = {
        "Content-Type": "application/json",
        "api-key": api_key,
    }
    return url, headers, request_body


@app.post("/api/llm/chat")
async def proxy_llm_chat(payload: ChatRequest, request: Request, session: Session = Depends(get_session)):
    url, headers, request_body = _build_llm_request(payload)

    # Log the conversation request (messages without exposing full content for privacy)
    message_summary = [{"role": m.role.value, "content_length": len(m.content)} for m in payload.messages]

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            response = await client.post(url, headers=headers, json=request_body)
    except httpx.HTTPError as exc:  # pragma: no cover - network safeguard
        log_action(session, "llm_chat_error", "llm", None, {"model": payload.model, "error": str(exc), "message_count": len(payload.messages)}, request)
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Upstream request failed: {exc}",
        ) from exc

    if response.status_code >= 400:
        log_action(session, "llm_chat_error", "llm", None, {"model": payload.model, "status_code": response.status_code, "message_count": len(payload.messages)}, request)
        raise HTTPException(
            status_code=response.status_code,
            detail=response.text,
        )

    data = response.json()
    choices = data.get("choices") or []
    if not choices:
        log_action(session, "llm_chat_error", "llm", None, {"model": payload.model, "error": "Empty choices in response"}, request)
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail="LLM returned empty response")
    message_content = choices[0].get("message", {}).get("content", "")

    # Extract thinking and text content from OpenAI's extended thinking response format
    # Content can be a string or an array of content blocks with types "thinking" and "text"
    thinking = None
    if isinstance(message_content, list):
        # Extract thinking blocks and text blocks separately
        thinking_parts = []
        text_parts = []
        for block in message_content:
            if isinstance(block, dict):
                if block.get("type") == "thinking":
                    thinking_parts.append(block.get("thinking", ""))
                elif block.get("type") == "text":
                    text_parts.append(block.get("text", ""))
        thinking = "\n".join(thinking_parts) if thinking_parts else None
        content = "\n".join(text_parts) if text_parts else ""
    else:
        content = message_content if message_content else ""

    # Log successful LLM conversation with full messages and response for auditing
    import json
    conversation_log = {
        "model": payload.model,
        "messages": [{"role": m.role.value, "content": m.content} for m in payload.messages],
        "response": content,
        "thinking": thinking,
        "usage": data.get("usage", {})
    }
    log_action(session, "llm_chat", "llm", None, conversation_log, request)

    return {"content": content, "thinking": thinking, "raw": data}


@app.post("/api/llm/chat/stream")
async def stream_llm_chat(payload: ChatRequest, request: Request, session: Session = Depends(get_session)):
    """
    Streaming LLM chat endpoint using Server-Sent Events (SSE).
    Streams tokens as they arrive, including tool calls.
    """
    import json as json_module

    url, headers, request_body = _build_llm_request(payload, stream=True)

    async def generate_events():
        accumulated_content = ""
        accumulated_tool_calls = []
        current_tool_call = None
        finish_reason = None

        try:
            async with httpx.AsyncClient(timeout=120) as client:
                async with client.stream("POST", url, headers=headers, json=request_body) as response:
                    if response.status_code >= 400:
                        error_text = ""
                        async for chunk in response.aiter_text():
                            error_text += chunk
                        yield f"data: {json_module.dumps({'error': error_text, 'status': response.status_code})}\n\n"
                        return

                    async for line in response.aiter_lines():
                        if not line:
                            continue

                        if line.startswith("data: "):
                            data_str = line[6:]

                            if data_str.strip() == "[DONE]":
                                # Send final message with complete data
                                final_event = {
                                    "type": "done",
                                    "content": accumulated_content,
                                    "tool_calls": accumulated_tool_calls if accumulated_tool_calls else None,
                                    "finish_reason": finish_reason,
                                }
                                yield f"data: {json_module.dumps(final_event)}\n\n"
                                break

                            try:
                                chunk_data = json_module.loads(data_str)
                                choices = chunk_data.get("choices") or []
                                if not choices:
                                    continue  # Skip chunks without choices
                                choice = choices[0]
                                delta = choice.get("delta", {})
                                finish_reason = choice.get("finish_reason") or finish_reason

                                # Handle content streaming
                                if "content" in delta and delta["content"]:
                                    content_chunk = delta["content"]
                                    accumulated_content += content_chunk
                                    yield f"data: {json_module.dumps({'type': 'content', 'content': content_chunk})}\n\n"

                                # Handle tool call streaming
                                if "tool_calls" in delta:
                                    for tc_delta in delta["tool_calls"]:
                                        tc_index = tc_delta.get("index", 0)

                                        # Ensure we have enough tool calls
                                        while len(accumulated_tool_calls) <= tc_index:
                                            accumulated_tool_calls.append({
                                                "id": "",
                                                "type": "function",
                                                "function": {"name": "", "arguments": ""}
                                            })

                                        current_tc = accumulated_tool_calls[tc_index]

                                        if "id" in tc_delta:
                                            current_tc["id"] = tc_delta["id"]

                                        if "function" in tc_delta:
                                            func_delta = tc_delta["function"]
                                            if "name" in func_delta:
                                                current_tc["function"]["name"] = func_delta["name"]
                                                # Emit tool call start event
                                                yield f"data: {json_module.dumps({'type': 'tool_call_start', 'index': tc_index, 'id': current_tc['id'], 'name': func_delta['name']})}\n\n"
                                            if "arguments" in func_delta:
                                                current_tc["function"]["arguments"] += func_delta["arguments"]

                            except json_module.JSONDecodeError:
                                continue

        except httpx.HTTPError as exc:
            yield f"data: {json_module.dumps({'type': 'error', 'error': str(exc)})}\n\n"
        except asyncio.CancelledError:
            # Client disconnected - gracefully end the stream
            yield f"data: {json_module.dumps({'type': 'error', 'error': 'Request cancelled'})}\n\n"
        except Exception as exc:
            # Catch all other exceptions to prevent TaskGroup errors
            yield f"data: {json_module.dumps({'type': 'error', 'error': str(exc)})}\n\n"

    return StreamingResponse(
        generate_events(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        }
    )


# PowerPoint Export Models and Endpoint
class SlideTaskItem(BaseModel):
    title: str
    date: str


class SlideUpdateItem(BaseModel):
    author: str
    date: str
    note: str


class SlideStakeholder(BaseModel):
    name: str
    team: str = ""


class SlideData(BaseModel):
    name: str
    description: str = ""
    executiveUpdate: str = ""
    targetDate: Optional[str] = None
    priority: str = "medium"
    status: str = "active"
    stakeholders: List[SlideStakeholder] = PydanticField(default_factory=list)
    recentlyCompleted: List[SlideTaskItem] = PydanticField(default_factory=list)
    nextUp: List[SlideTaskItem] = PydanticField(default_factory=list)
    recentUpdates: List[SlideUpdateItem] = PydanticField(default_factory=list)


class SlidesExportPayload(BaseModel):
    slides: List[SlideData]


def create_powerpoint_presentation(slides: List[SlideData]) -> bytes:
    """Generate a PowerPoint presentation from slide data matching the web UI design."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 standard width
    prs.slide_height = Inches(7.5)    # 16:9 standard height

    # Define colors matching the web UI (CSS variables)
    CHARCOAL = RGBColor(58, 54, 49)      # --charcoal: #3a3631
    STONE = RGBColor(107, 101, 84)       # --stone: #6b6554
    CLOUD = RGBColor(232, 227, 216)      # --cloud: #e8e3d8
    CREAM = RGBColor(250, 248, 243)      # --cream: #faf8f3
    SAGE = RGBColor(122, 155, 118)       # --sage: #7a9b76
    CORAL = RGBColor(215, 119, 100)      # --coral: #d77764
    AMBER = RGBColor(218, 165, 32)       # --amber: #daa520
    EARTH = RGBColor(139, 111, 71)       # --earth: #8b6f47
    WHITE = RGBColor(255, 255, 255)

    def sanitize_text(value: str) -> str:
        """Remove control characters that can corrupt the PPTX XML."""
        if not value:
            return ""
        return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", value)

    def get_priority_color(priority: str) -> RGBColor:
        colors = {
            'high': CORAL,
            'medium': AMBER,
            'low': SAGE
        }
        return colors.get(priority, STONE)

    def add_rounded_rectangle(slide, left, top, width, height, fill_color=None, line_color=None):
        """Add a rounded rectangle shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color if fill_color else WHITE
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        # Set corner radius - use try/except to handle shapes without adjustments
        try:
            if shape.adjustments and len(shape.adjustments) > 0:
                shape.adjustments[0] = 0.1
        except (IndexError, TypeError):
            pass  # Shape doesn't support adjustments
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=12, font_color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT):
        """Add a text box with specified formatting."""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = sanitize_text(text)
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Arial"
        p.alignment = alignment
        return textbox

    def format_datetime_simple(date_str: str) -> str:
        """Format datetime string for display."""
        try:
            date = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
            return date.strftime('%b %d, %Y')
        except:
            return date_str

    # Create a slide for each project
    for slide_data in slides:
        # Add a blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add background - using slide background property instead of shape
        # to avoid z-order issues and file corruption
        slide_background = slide.background
        fill = slide_background.fill
        fill.solid()
        fill.fore_color.rgb = CREAM

        # Layout constants
        MARGIN = Inches(0.5)
        HEADER_HEIGHT = Inches(1.2)
        PANEL_GAP = Inches(0.2)
        CONTENT_TOP = MARGIN + HEADER_HEIGHT + Inches(0.3)
        CONTENT_HEIGHT = prs.slide_height - CONTENT_TOP - MARGIN
        HALF_WIDTH = (prs.slide_width - MARGIN * 2 - PANEL_GAP) / 2

        # ===== HEADER SECTION =====
        # Project name
        title_box = add_text_box(
            slide, MARGIN, MARGIN, prs.slide_width - MARGIN * 2, Inches(0.5),
            sanitize_text(slide_data.name), font_size=24, font_color=CHARCOAL, bold=True
        )

        # Project description (subtitle)
        if slide_data.description:
            add_text_box(
                slide, MARGIN, MARGIN + Inches(0.45), prs.slide_width - MARGIN * 2, Inches(0.3),
                sanitize_text(slide_data.description), font_size=12, font_color=CHARCOAL
            )

        # Metadata row
        meta_y = MARGIN + Inches(0.85)
        meta_x = MARGIN

        # Target date
        target_text = sanitize_text(f"Target: {slide_data.targetDate if slide_data.targetDate else 'TBD'}")
        target_box = add_text_box(slide, meta_x, meta_y, Inches(1.5), Inches(0.25), target_text, font_size=10, font_color=STONE)
        meta_x += Inches(1.6)

        # Priority badge
        priority_color = get_priority_color(slide_data.priority)
        priority_shape = add_rounded_rectangle(slide, meta_x, meta_y, Inches(1.1), Inches(0.25), fill_color=CREAM, line_color=priority_color)
        priority_tf = priority_shape.text_frame
        priority_tf.paragraphs[0].text = sanitize_text(f"{slide_data.priority} priority")
        priority_tf.paragraphs[0].font.size = Pt(9)
        priority_tf.paragraphs[0].font.color.rgb = priority_color
        priority_tf.paragraphs[0].font.bold = True
        priority_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        priority_shape.text_frame.paragraphs[0].space_before = Pt(2)
        meta_x += Inches(1.2)

        # Status badge
        status_shape = add_rounded_rectangle(slide, meta_x, meta_y, Inches(0.9), Inches(0.25), fill_color=CLOUD)
        status_tf = status_shape.text_frame
        status_tf.paragraphs[0].text = sanitize_text(slide_data.status)
        status_tf.paragraphs[0].font.size = Pt(9)
        status_tf.paragraphs[0].font.color.rgb = CHARCOAL
        status_tf.paragraphs[0].font.bold = True
        status_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        status_shape.text_frame.paragraphs[0].space_before = Pt(2)
        meta_x += Inches(1.0)

        # Stakeholders
        if slide_data.stakeholders:
            stakeholder_names = [s.name for s in slide_data.stakeholders[:5]]
            stakeholder_text = sanitize_text(", ".join(stakeholder_names))
            if len(slide_data.stakeholders) > 5:
                stakeholder_text += sanitize_text(f" +{len(slide_data.stakeholders) - 5}")
            add_text_box(slide, meta_x, meta_y, Inches(4), Inches(0.25), sanitize_text(f"Team: {stakeholder_text}"), font_size=10, font_color=STONE)

        # Header divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, CONTENT_TOP - Inches(0.15), prs.slide_width - MARGIN * 2, Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = CLOUD
        line.line.fill.background()

        # ===== LEFT COLUMN =====
        left_x = MARGIN
        panel_width = HALF_WIDTH
        exec_height = Inches(1.5)
        updates_height = CONTENT_HEIGHT - exec_height - PANEL_GAP

        # Executive Update Panel
        exec_panel = add_rounded_rectangle(slide, left_x, CONTENT_TOP, panel_width, exec_height, fill_color=WHITE, line_color=CLOUD)

        # Executive Update title
        add_text_box(slide, left_x + Inches(0.15), CONTENT_TOP + Inches(0.1), panel_width - Inches(0.3), Inches(0.2),
                    "EXECUTIVE UPDATE", font_size=9, font_color=STONE, bold=True)

        # Executive Update content
        exec_content = sanitize_text(slide_data.executiveUpdate or slide_data.description or "No executive update yet.")
        exec_text_box = slide.shapes.add_textbox(
            left_x + Inches(0.15), CONTENT_TOP + Inches(0.35),
            panel_width - Inches(0.3), exec_height - Inches(0.5)
        )
        exec_tf = exec_text_box.text_frame
        exec_tf.word_wrap = True
        exec_tf.paragraphs[0].text = exec_content[:500]  # Limit text length
        exec_tf.paragraphs[0].font.size = Pt(10)
        exec_tf.paragraphs[0].font.color.rgb = STONE
        exec_tf.paragraphs[0].font.name = "Arial"

        # Recent Updates Panel
        updates_top = CONTENT_TOP + exec_height + PANEL_GAP
        updates_panel = add_rounded_rectangle(slide, left_x, updates_top, panel_width, updates_height, fill_color=WHITE, line_color=CLOUD)

        # Recent Updates title
        add_text_box(slide, left_x + Inches(0.15), updates_top + Inches(0.1), panel_width - Inches(0.3), Inches(0.2),
                    "RECENT UPDATES", font_size=9, font_color=STONE, bold=True)

        # Recent Updates content
        updates_y = updates_top + Inches(0.35)
        for i, update in enumerate(slide_data.recentUpdates[:3]):
            if updates_y > updates_top + updates_height - Inches(0.4):
                break
            # Author and date
            add_text_box(slide, left_x + Inches(0.15), updates_y, Inches(1.5), Inches(0.2),
                        sanitize_text(update.author), font_size=10, font_color=CHARCOAL, bold=True)
            add_text_box(slide, left_x + panel_width - Inches(1.5), updates_y, Inches(1.35), Inches(0.2),
                        sanitize_text(format_datetime_simple(update.date)), font_size=9, font_color=STONE, alignment=PP_ALIGN.RIGHT)
            updates_y += Inches(0.2)
            # Note text
            note_box = slide.shapes.add_textbox(left_x + Inches(0.15), updates_y, panel_width - Inches(0.3), Inches(0.4))
            note_tf = note_box.text_frame
            note_tf.word_wrap = True
            note_tf.paragraphs[0].text = sanitize_text(update.note)[:150]
            note_tf.paragraphs[0].font.size = Pt(9)
            note_tf.paragraphs[0].font.color.rgb = STONE
            note_tf.paragraphs[0].font.name = "Arial"
            updates_y += Inches(0.45)

        if not slide_data.recentUpdates:
            add_text_box(slide, left_x + Inches(0.15), updates_top + Inches(0.4), panel_width - Inches(0.3), Inches(0.2),
                        "No updates yet.", font_size=10, font_color=STONE)

        # ===== RIGHT COLUMN =====
        right_x = MARGIN + HALF_WIDTH + PANEL_GAP
        half_height = (CONTENT_HEIGHT - PANEL_GAP) / 2

        # Recently Completed Panel
        completed_panel = add_rounded_rectangle(slide, right_x, CONTENT_TOP, panel_width, half_height, fill_color=WHITE, line_color=SAGE)

        # Recently Completed title
        add_text_box(slide, right_x + Inches(0.15), CONTENT_TOP + Inches(0.1), panel_width - Inches(0.3), Inches(0.2),
                    "RECENTLY COMPLETED", font_size=9, font_color=STONE, bold=True)

        # Recently Completed content
        completed_y = CONTENT_TOP + Inches(0.35)
        for task in slide_data.recentlyCompleted[:3]:
            if completed_y > CONTENT_TOP + half_height - Inches(0.3):
                break
            # Task title
            task_box = slide.shapes.add_textbox(right_x + Inches(0.15), completed_y, panel_width - Inches(1.2), Inches(0.25))
            task_tf = task_box.text_frame
            task_tf.word_wrap = True
            task_tf.paragraphs[0].text = sanitize_text(task.title)[:60]
            task_tf.paragraphs[0].font.size = Pt(10)
            task_tf.paragraphs[0].font.color.rgb = CHARCOAL
            task_tf.paragraphs[0].font.bold = True
            task_tf.paragraphs[0].font.name = "Arial"
            # Date
            add_text_box(slide, right_x + panel_width - Inches(1), completed_y, Inches(0.85), Inches(0.2),
                        sanitize_text(task.date), font_size=9, font_color=SAGE, alignment=PP_ALIGN.RIGHT)
            completed_y += Inches(0.35)

        if not slide_data.recentlyCompleted:
            add_text_box(slide, right_x + Inches(0.15), CONTENT_TOP + Inches(0.4), panel_width - Inches(0.3), Inches(0.2),
                        "No recently completed tasks.", font_size=10, font_color=STONE)

        # Next Up Panel
        nextup_top = CONTENT_TOP + half_height + PANEL_GAP
        nextup_panel = add_rounded_rectangle(slide, right_x, nextup_top, panel_width, half_height, fill_color=WHITE, line_color=SAGE)

        # Next Up title
        add_text_box(slide, right_x + Inches(0.15), nextup_top + Inches(0.1), panel_width - Inches(0.3), Inches(0.2),
                    "NEXT UP", font_size=9, font_color=STONE, bold=True)

        # Next Up content
        nextup_y = nextup_top + Inches(0.35)
        for task in slide_data.nextUp[:3]:
            if nextup_y > nextup_top + half_height - Inches(0.3):
                break
            # Task title
            task_box = slide.shapes.add_textbox(right_x + Inches(0.15), nextup_y, panel_width - Inches(1.2), Inches(0.25))
            task_tf = task_box.text_frame
            task_tf.word_wrap = True
            task_tf.paragraphs[0].text = sanitize_text(task.title)[:60]
            task_tf.paragraphs[0].font.size = Pt(10)
            task_tf.paragraphs[0].font.color.rgb = CHARCOAL
            task_tf.paragraphs[0].font.bold = True
            task_tf.paragraphs[0].font.name = "Arial"
            # Date - color based on content (overdue = coral, otherwise stone)
            date_text = sanitize_text(task.date)
            date_color = CORAL if 'overdue' in date_text.lower() else (AMBER if 'today' in date_text.lower() or 'tomorrow' in date_text.lower() else STONE)
            add_text_box(slide, right_x + panel_width - Inches(1), nextup_y, Inches(0.85), Inches(0.2),
                        date_text, font_size=9, font_color=date_color, alignment=PP_ALIGN.RIGHT)
            nextup_y += Inches(0.35)

        if not slide_data.nextUp:
            add_text_box(slide, right_x + Inches(0.15), nextup_top + Inches(0.4), panel_width - Inches(0.3), Inches(0.2),
                        "No upcoming tasks.", font_size=10, font_color=STONE)

    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()


@app.post("/api/slides/export")
def export_slides_to_powerpoint(payload: SlidesExportPayload, request: Request, session: Session = Depends(get_session)):
    """Export slides to PowerPoint format."""
    if not payload.slides:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No slides provided")

    try:
        pptx_bytes = create_powerpoint_presentation(payload.slides)
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="python-pptx library not installed. Please install it with: pip install python-pptx"
        )
    except Exception as e:
        logger.exception("Failed to generate PowerPoint")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate PowerPoint: {str(e)}"
        )

    try:
        from pptx import Presentation

        Presentation(io.BytesIO(pptx_bytes))
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="python-pptx library not installed. Please install it with: pip install python-pptx"
        )
    except Exception:
        logger.exception("Generated PowerPoint failed validation")
        return JSONResponse(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            content={"detail": "Generated PowerPoint file is invalid and could not be read."}
        )

    log_action(session, "export_slides_pptx", "slides", None, {"slide_count": len(payload.slides)}, request)

    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=portfolio-slides-{datetime.utcnow().strftime('%Y-%m-%d')}.pptx"
        }
    )


@app.get("/")
def root():
    return {"status": "ok"}


def validate_admin_token_for_cli(admin_token: str | None) -> None:
    expected_token = os.getenv(ADMIN_TOKEN_ENV)
    if not expected_token:
        if current_environment() in PROTECTED_ENVIRONMENTS:
            raise RuntimeError("Admin token not configured for this environment")
        logger.warning(
            "Admin token not configured; allowing admin CLI command in non-protected environment"
        )
        return
    if not admin_token:
        raise RuntimeError("Admin token required for this operation")
    if admin_token != expected_token:
        raise RuntimeError("Invalid admin token")


def run_people_backfill_cli(admin_token: str | None) -> None:
    validate_admin_token_for_cli(admin_token)
    create_db_and_tables()
    with Session(engine) as session:
        run_people_backfill(session)
    logger.info("People backfill completed or already applied.")


def run_unique_constraints_migration_cli(admin_token: str | None) -> None:
    validate_admin_token_for_cli(admin_token)
    create_db_and_tables()
    with Session(engine) as session:
        migrate_add_unique_constraints(session)
    logger.info("UNIQUE constraints migration completed or already applied.")


def main() -> None:
    parser = argparse.ArgumentParser(description="Manity Portfolio API utilities")
    subparsers = parser.add_subparsers(dest="command")

    backfill_parser = subparsers.add_parser(
        "run-people-backfill",
        help="Run the people backfill migration before starting the API server",
    )
    backfill_parser.add_argument(
        "--admin-token",
        default=os.getenv(ADMIN_TOKEN_ENV),
        help="Admin token for protected environments (defaults to MANITY_ADMIN_TOKEN)",
    )

    migration_parser = subparsers.add_parser(
        "run-unique-constraints-migration",
        help="Add UNIQUE constraints to project.name and person.name",
    )
    migration_parser.add_argument(
        "--admin-token",
        default=os.getenv(ADMIN_TOKEN_ENV),
        help="Admin token for protected environments (defaults to MANITY_ADMIN_TOKEN)",
    )

    args = parser.parse_args()

    if args.command == "run-people-backfill":
        run_people_backfill_cli(args.admin_token)
    elif args.command == "run-unique-constraints-migration":
        run_unique_constraints_migration_cli(args.admin_token)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
